#!/usr/bin/env python3
"""
kpi_renderer.py — canonical KPI rendering для build_v9.

Зачем: donor 43/44 в шаблоне = guide-слайды (показывают примеры размеров),
не доноры для произвольных KPI. Их placeholders расположены обучающе
(199pt + 199pt + 44pt) — не для real data.

Решение: рисуем KPI shapes с нуля на blank donor с CLEAN GRID:
- 1 number → hero center
- 2 numbers → 50/50 grid
- 3 numbers → 33/33/33 grid (все одного размера)
- Title 32pt SemiBold top
- Numbers 130pt SemiBold (помещается без overflow)
- Descriptions 16pt Regular under numbers
- Accent: 1 of N number = #26D07C (canonical §12)

Layout (1280×720 slide):
  Title:           (35, 38, 1209, 53) px       32pt SemiBold #222222
  Number row:      y=220, height=280
    1 number:      x=440, w=400 (centered)
    2 numbers:     x=140, w=400 + x=740, w=400
    3 numbers:     x=35, w=400 + x=440, w=400 + x=840, w=400
  Desc row:        y=530, height=100
    Same x columns
  Optional %:      next to number, top-right corner of number box
"""
import os
import sys
import json
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from text_sanitize import sanitize_text


def _load_template_version():
    """Грузит brand/template-version.json (маппинг slide-индексов). Fallback — None."""
    here = os.path.dirname(os.path.abspath(__file__))
    candidates = [
        os.path.join(here, "..", "brand", "template-version.json"),
        os.path.join(os.getcwd(), "brand", "template-version.json"),
        os.path.join(os.getcwd(), "pptx-skill", "brand", "template-version.json"),
    ]
    for path in candidates:
        try:
            if os.path.isfile(path):
                with open(path, encoding="utf-8") as f:
                    return json.load(f)
        except Exception:
            continue
    return None


_TPL_VER = _load_template_version()


GRAPHITE = RGBColor(0x22, 0x22, 0x22)
GREEN = RGBColor(0x26, 0xD0, 0x7C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "SB Sans Display"
# Полужирное = отдельный font face (встроен в шаблон), НЕ bold-флаг.
# Canonical 2026-05-29 (Problem #3): Bold запрещён — эмфаза только через SemiBold.
FONT_SEMIBOLD = "SB Sans Display Semibold"


# Canonical §4 (slide 43): фрейм 199pt помещает максимум 2 значащих цифры.
KPI_HERO_DIGIT_LIMIT = 2


def _count_significant_digits(value):
    """Считает значащие цифры в KPI value, игнорируя разделители и unit-suffix.

    Примеры:
      '12'    → 2 цифры
      '99%'   → 2 (% игнорируется — pct идёт отдельным флагом)
      '1.5'   → 2 (точка игнорируется, но 1+5=2)
      '1.5K'  → 2 (K — буква-сокращение, не считается)
      '150'   → 3 (overflow!)
      '~150'  → 3 (тильда игнорируется)
      '12,5'  → 3 (русская десятичная — 3 цифры в frame не помещаются)
    """
    if value is None:
        return 0
    s = str(value)
    digits = "".join(c for c in s if c.isdigit())
    return len(digits)


def _add_text_box(slide, left_px, top_px, width_px, height_px, text,
                  font_size_pt=14, font_name=None,
                  bold=False, color=GRAPHITE, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP):
    """Add text box with canonical font/size/color.

    Эмфаза (bold=True) реализуется через начертание SemiBold, а не bold-флаг
    (Problem #3 2026-05-29). font_name=None → авто: Regular или Semibold по bold.
    """
    EMU = 9525
    box = slide.shapes.add_textbox(
        Emu(left_px * EMU), Emu(top_px * EMU),
        Emu(width_px * EMU), Emu(height_px * EMU)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    # kpi_native slides are skipped by apply_kpi_emphasis (render_kpi styles
    # numbers itself), so these texts never feed the ** emphasis path — safe to
    # strip both markdown and control chars.
    run.text = sanitize_text(text)
    # SemiBold через face, Bold-флаг не используем
    if font_name is None:
        run.font.name = FONT_SEMIBOLD if bold else FONT
    else:
        run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = False
    run.font.color.rgb = color
    return box


def _add_accent_bar(slide, left_px, top_px, width_px, height_px, color=GREEN):
    """Зелёная подчёркивающая плашка-акцент под главной KPI-цифрой.

    Canonical 2026-05-29 (Problem #2): акцент НЕ цветом текста, а отдельным
    цветным элементом. Сама цифра остаётся #222222 (светлый) / белой (тёмный),
    а «главный показатель» помечается этой зелёной плашкой.
    """
    EMU = 9525
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(left_px * EMU), Emu(top_px * EMU),
        Emu(width_px * EMU), Emu(height_px * EMU),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # без рамки
    # Снять любые эффекты — и явные, и тематические (effectRef). Плоско, без исключений.
    from effects_util import strip_effects
    strip_effects(shape._element)
    return shape


def _is_title_placeholder(sp):
    """sp (<p:sp>) — это title/ctrTitle placeholder шаблона?"""
    nv = sp.find(qn("p:nvSpPr"))
    if nv is None:
        return False
    nvpr = nv.find(qn("p:nvPr"))
    if nvpr is None:
        return False
    ph = nvpr.find(qn("p:ph"))
    if ph is None:
        return False
    return ph.get("type") in ("title", "ctrTitle")


# D4 (2026-06-06): the active content zone (px → EMU). A decorative pic that
# intersects this rectangle would collide with rendered content, so it is
# removed even when keep_decor is on; pics fully outside (corners/edges) stay.
_EMU_PER_PX = 9525
CONTENT_ZONE_PX = (35, 120, 1245, 660)  # left, top, right, bottom
CONTENT_ZONE_EMU = tuple(v * _EMU_PER_PX for v in CONTENT_ZONE_PX)


def _pic_overlaps_content(child) -> bool:
    """True if a <p:pic> element's bbox intersects CONTENT_ZONE_EMU. Missing
    geometry → treat as overlapping (conservative: strip it)."""
    spPr = child.find(qn("p:spPr"))
    if spPr is None:
        return True
    xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is None:
        return True
    off = xfrm.find(qn("a:off"))
    ext = xfrm.find(qn("a:ext"))
    if off is None or ext is None:
        return True
    try:
        x = int(off.get("x")); y = int(off.get("y"))
        cx = int(ext.get("cx")); cy = int(ext.get("cy"))
    except (TypeError, ValueError):
        return True
    l, t, r, b = CONTENT_ZONE_EMU
    return not (x + cx <= l or x >= r or y + cy <= t or y >= b)


def clean_slide_to_blank(slide, keep_title=True, keep_decor=False):
    """Удалить все shapes на slide, КРОМЕ title-placeholder шаблона (если
    keep_title). При keep_decor=True декоративные <p:pic> вне контент-зоны
    сохраняются (D4 reclaim), а пересекающие контент-зону — удаляются.
    Layout-inherited (logo, footer) останутся."""
    spTree = slide.shapes._spTree
    to_remove = []
    for child in list(spTree):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('sp', 'pic', 'graphicFrame', 'cxnSp', 'grpSp'):
            if keep_title and tag == 'sp' and _is_title_placeholder(child):
                # сохранить placeholder, но очистить донорский текст
                txBody = child.find(qn('p:txBody'))
                if txBody is not None:
                    for p_el in txBody.findall(qn('a:p')):
                        for r_el in p_el.findall(qn('a:r')):
                            p_el.remove(r_el)
                continue
            if keep_decor and tag == 'pic' and not _pic_overlaps_content(child):
                continue
            to_remove.append(child)
    for el in to_remove:
        spTree.remove(el)


def set_slide_title(slide, text, dark=False):
    """Вписать заголовок слайда в штатный TITLE-placeholder шаблона.

    Canonical (Problem #6, 2026-05-29): единая позиция/размер из шаблона
    (35,38 / 963×54 / 20pt SemiBold CAPS, графит/белый). Если placeholder на
    слайде нет — fallback на канонический textbox той же геометрии.
    """
    if not text:
        return None
    color = WHITE if dark else GRAPHITE
    title_ph = None
    try:
        title_ph = slide.shapes.title
    except Exception:
        title_ph = None
    if title_ph is None:
        return _add_text_box(slide, 35, 38, 963, 54, text.upper(),
                             font_size_pt=20, bold=True, color=color,
                             anchor=MSO_ANCHOR.MIDDLE)
    tf = title_ph.text_frame
    tf.clear()
    run = tf.paragraphs[0].add_run()
    # Generic title setter (chart/table/flow). Titles never carry intentional
    # ** emphasis and are not on the kpi_emphasis path, so strip both markdown
    # and control chars (fixes _X000B_ / ** leak in native-slide titles).
    run.text = sanitize_text(text).upper()
    run.font.name = FONT_SEMIBOLD
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = color
    return title_ph


def render_kpi(slide, kpi_config, dark=False):
    """
    Render KPI block on slide.

    kpi_config = {
        "title": "Что получилось",
        "numbers": [
            {"value": "12", "desc": "сократили AI-инициатив", "accent": False, "pct": False},
            {"value": "84", "desc": "вовлечённость", "accent": False, "pct": True},
            {"value": "5",  "desc": "инструментов", "accent": True, "pct": False}
        ]
    }

    dark: True if slide is dark (text white instead of graphite)
    """
    text_color = WHITE if dark else GRAPHITE
    n = len(kpi_config.get("numbers", []))

    if n == 0 or n > 3:
        raise ValueError(f"KPI supports 1-3 numbers, got {n}")

    # Title — в штатный placeholder шаблона (Problem #6)
    title = kpi_config.get("title", "")
    if title:
        set_slide_title(slide, title, dark=dark)

    # D3 (2026-06-06): inline layout — number LEFT, description RIGHT, vertically
    # centered (template slide 43). Color stays graphite (canon 2026-05-29).
    if n == 1:
        col_x = [60]
        col_w = [1160]
    elif n == 2:
        col_x = [60, 660]
        col_w = [560, 560]
    else:  # n == 3
        col_x = [40, 440, 840]
        col_w = [380, 380, 380]

    ROW_TOP = 240
    ROW_HEIGHT = 240
    NUMBER_FONT = 150 if n == 1 else (110 if n == 2 else 90)
    if n == 3:
        max_chars = max(len(str(x["value"])) for x in kpi_config["numbers"])
        if max_chars > 3:
            NUMBER_FONT = 76
    DESC_FONT = 16 if n == 1 else (14 if n == 2 else 12)

    for i, num in enumerate(kpi_config["numbers"]):
        x = col_x[i]
        w = col_w[i]
        color = text_color  # graphite/white — never the accent itself
        value = num["value"]
        has_pct = num.get("pct", False)

        # Defensive normalize: value может уже нести хвостовой '%' (LLM
        # проигнорировал контракт «value без %, pct — флаг»). Срезаем его и
        # выводим через pct-надстрочник, иначе '%' рендерится дважды и
        # переполняет фрейм (live run 2026-06-05, slide7 '99.97%').
        if isinstance(value, str) and value.rstrip().endswith("%"):
            value = value.rstrip()[:-1].rstrip()
            has_pct = True

        # Canonical §4: крупный hero-frame ⇒ max 2 значащих цифры, иначе overflow.
        if NUMBER_FONT >= 150:
            digits = _count_significant_digits(value)
            if digits > KPI_HERO_DIGIT_LIMIT:
                print(
                    f"WARN: KPI value '{value}' содержит {digits} значащих цифр "
                    f"при {NUMBER_FONT}pt frame (canonical max {KPI_HERO_DIGIT_LIMIT}).",
                    file=sys.stderr,
                )

        # Number box: left-aligned, width sized to the value (~0.62×font per
        # digit). Shrink the font so the value always fits on ONE line within
        # the column (≥120px reserved for the description), then disable wrap —
        # otherwise long values like "99.97" break to 2 lines (live deck_a.s7).
        DESC_MIN = 120
        avail_num_w = max(160, w - DESC_MIN)
        nchars = max(1, len(str(value)))
        nf = NUMBER_FONT
        nat_w = int(nchars * nf * 0.62) + 40
        if nat_w > avail_num_w:
            nf = max(40, int((avail_num_w - 40) / (nchars * 0.62)))
            nat_w = int(nchars * nf * 0.62) + 40
        num_w = min(max(120, nat_w), avail_num_w)
        num_box = _add_text_box(slide, x, ROW_TOP, num_w, ROW_HEIGHT,
                      value, font_size_pt=nf, bold=False, color=color,
                      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        num_box.text_frame.word_wrap = False

        # Enlarged %: ≈0.5× number height, kerned to the number's top-right.
        if has_pct:
            pct_size = max(40, int(nf * 0.55))
            pct_x = x + num_w - 16
            pct_y = ROW_TOP + 10
            _add_text_box(slide, pct_x, pct_y, 90, 120, "%",
                          font_size_pt=pct_size, bold=True, color=color,
                          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

        # Green underline-plate accent under the number (canon: accent is the
        # plate, not a green number).
        if num.get("accent", False):
            _add_accent_bar(slide, x, ROW_TOP + ROW_HEIGHT - 8,
                            num_w, 8, color=GREEN)

        # Description: to the RIGHT of the number, vertically centered to it.
        desc = num.get("desc", "")
        if desc:
            desc_x = x + num_w + 16
            desc_w = x + w - desc_x
            _add_text_box(slide, desc_x, ROW_TOP, max(80, desc_w), ROW_HEIGHT,
                          desc, font_size_pt=DESC_FONT, bold=False,
                          color=text_color,
                          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


# Constants for blank donor selection.
# Источник: brand/template-version.json → blank_donors (white/dark).
# Hardcoded fallback — на случай отсутствия файла.
# dark=51 (layout 'Контент / темный 1', bg #222222); было 22 (белый layout — баг).
_BLANK_DEFAULTS = {"white": 30, "dark": 51}
_blank_cfg = (_TPL_VER or {}).get("blank_donors", {}) if _TPL_VER else {}
BLANK_DONOR_WHITE = int(_blank_cfg.get("white", _BLANK_DEFAULTS["white"]))
BLANK_DONOR_DARK = int(_blank_cfg.get("dark", _BLANK_DEFAULTS["dark"]))


# Standalone test
if __name__ == "__main__":
    from pptx import Presentation
    from template_path import resolve_template

    p = Presentation(resolve_template())
    # Use slide 30 as clean blank donor — but FULLY clean shapes
    slide = list(p.slides)[BLANK_DONOR_WHITE - 1]
    clean_slide_to_blank(slide)

    # Render KPI
    render_kpi(slide, {
        "title": "ЧТО ПОЛУЧИЛОСЬ ЧЕРЕЗ 6 МЕСЯЦЕВ",
        "numbers": [
            {"value": "12", "desc": "сократили AI-инициатив", "pct": True},
            {"value": "84", "desc": "вовлечённость команды", "pct": True},
            {"value": "5",  "desc": "рабочих инструментов", "accent": True}
        ]
    }, dark=False)

    out = "pptx-skill/output/kpi_renderer_test.pptx"
    p.save(out)
    print(f"Saved {out}, used slide 30 as blank donor")
