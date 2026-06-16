#!/usr/bin/env python3
"""
chart_native_pptx.py — редактируемая (native) PowerPoint chart через python-pptx.

В отличие от chart_engine.py (который рендерит matplotlib → PNG → image), этот
модуль создаёт настоящий PPTX chart — пользователь в PowerPoint может «Edit Data»
и менять цифры через Excel-окно.

Поддерживаемые типы:
  area_stacked  → XL_CHART_TYPE.AREA_STACKED
  area_100      → XL_CHART_TYPE.AREA_STACKED_100
  bar           → XL_CHART_TYPE.COLUMN_CLUSTERED
  bar_stacked   → XL_CHART_TYPE.COLUMN_STACKED
  line          → XL_CHART_TYPE.LINE
  pie           → XL_CHART_TYPE.PIE

Canonical palette — из chart_engine (Cloud.ru pastel + GREEN accent).
"""
import os
import json
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

from text_sanitize import sanitize_text


def _hex_to_rgb_color(hx):
    h = hx.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _load_palette():
    here = os.path.dirname(os.path.abspath(__file__))
    candidates = [
        os.path.join(here, "..", "brand", "palette.json"),
        os.path.join(os.getcwd(), "brand", "palette.json"),
        os.path.join(os.getcwd(), "pptx-skill", "brand", "palette.json"),
    ]
    for path in candidates:
        try:
            if os.path.isfile(path):
                with open(path, encoding="utf-8") as f:
                    return json.load(f)
        except Exception:
            continue
    return None


_PAL = _load_palette()


def _color(name, fallback_hex):
    """Из brand/palette.json (base/base_alts) или fallback."""
    if _PAL:
        for section in ("base", "base_alts"):
            hx = _PAL.get(section, {}).get(name)
            if hx:
                return _hex_to_rgb_color(hx)
    return _hex_to_rgb_color(fallback_hex)


GRAPHITE = _color("Black", "#222222")
GREEN = _color("Green", "#26D07C")
WHITE = _color("White", "#FFFFFF")


# Не-accent серии: chart_extension из palette.json. GREEN зарезервирован под accent_idx.
# Порядок подобран для максимальной различимости: тёмный → синий → песочный → мята → серый.
def _build_non_accent_colors():
    if _PAL:
        ext = _PAL.get("chart_extension", {})
        # Желаемый порядок (по дизайну):
        order = ["dark_gray", "Blue", "pastel_yellow", "pastel_mint", "mid_gray", "pastel_green_tint"]
        result = []
        for key in order:
            hx = ext.get(key) or _PAL.get("base", {}).get(key)
            if hx:
                result.append(_hex_to_rgb_color(hx))
        if len(result) >= 5:
            return result
    return [
        _hex_to_rgb_color("#666666"),
        _hex_to_rgb_color("#C0E0FC"),
        _hex_to_rgb_color("#E8F5C7"),
        _hex_to_rgb_color("#C9F2EA"),
        _hex_to_rgb_color("#BDBDBD"),
        _hex_to_rgb_color("#A8E5C9"),
    ]


NON_ACCENT_COLORS = _build_non_accent_colors()


def _green_ramp(n: int, accent_idx: int = -1) -> list:
    """Return n RGBColors ramping dark→bright green; the accent bar gets the
    brand green (#26D07C, brightest). Used by the editorial single-series path
    where each bar is its own single-value series."""
    if n <= 0:
        return []
    # Dark→bright endpoints (graphite-green → brand green).
    lo = (0x0E, 0x5A, 0x38)
    hi = (0x26, 0xD0, 0x7C)
    out = []
    for i in range(n):
        t = i / max(n - 1, 1)
        rgb = tuple(int(lo[c] + (hi[c] - lo[c]) * t) for c in range(3))
        out.append(RGBColor(*rgb))
    if 0 <= accent_idx < n:
        out[accent_idx] = RGBColor(0x26, 0xD0, 0x7C)
    return out


def is_editorial_eligible(chart_config: dict) -> bool:
    """Editorial path applies only to single-series bar charts marked
    style='editorial'. Everything else keeps the clean uniform chart."""
    if chart_config.get("style") != "editorial":
        return False
    if chart_config.get("type") != "bar":
        return False
    series = chart_config.get("series") or []
    return len(series) == 1


CHART_TYPE_MAP = {
    "area_stacked": XL_CHART_TYPE.AREA_STACKED,
    "area_100": XL_CHART_TYPE.AREA_STACKED_100,
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
}


def _apply_series_colors(chart, accent_idx=-1):
    """Окрашивает серии по canonical palette. accent_idx получает GREEN,
    остальные распределяются по NON_ACCENT_COLORS в исходном порядке."""
    non_accent_iter = iter(NON_ACCENT_COLORS)
    for i, s in enumerate(chart.series):
        if i == accent_idx:
            color = GREEN
        else:
            try:
                color = next(non_accent_iter)
            except StopIteration:
                color = NON_ACCENT_COLORS[-1]
        fill = s.format.fill
        fill.solid()
        fill.fore_color.rgb = color
        line = s.format.line
        line.color.rgb = color
        line.width = Emu(0)


# Полужирное = отдельный font face (встроен в шаблон), НЕ bold-флаг.
# Canonical 2026-05-29 (Problem #3): Bold запрещён — эмфаза только через SemiBold.
FONT = "SB Sans Display"
FONT_SEMIBOLD = "SB Sans Display Semibold"


def _style_text(tf, color=GRAPHITE, size_pt=10, bold=False):
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.color.rgb = color
            r.font.size = Pt(size_pt)
            r.font.bold = False  # эмфаза через SemiBold-face, не bold-флаг
            r.font.name = FONT_SEMIBOLD if bold else FONT


def add_chart_to_slide(slide, chart_config, left, top, width, height, dark=False):
    """
    Добавляет редактируемую chart на slide.

    chart_config:
      type: area_stacked|area_100|bar|bar_stacked|line|pie
      title: optional строка
      x: list[str|int] — категории (для не-pie)
      labels: list[str] — для pie (вместо x)
      values: list[float] — для pie (вместо series)
      series: list[{"name": str, "data": list[float]}] — для не-pie
      accent_idx: int — индекс зелёного hero (default -1)
      x_label / y_label: optional строки
    """
    ctype = CHART_TYPE_MAP.get(chart_config.get("type", "bar"))
    if ctype is None:
        raise ValueError(f"Unsupported chart type: {chart_config.get('type')}")

    editorial = is_editorial_eligible(chart_config)
    accent_idx = chart_config.get("accent_idx", -1)

    cd = CategoryChartData()
    if ctype == XL_CHART_TYPE.PIE:
        # Canonical ChartConfig shape (parse_pptx, schemas/slides.py) carries
        # ``x`` + ``series`` for EVERY chart type. Legacy / LLM-native callers
        # may still pass ``labels`` + ``values`` — accept those as a fallback.
        pie_cats = chart_config.get("x")
        if pie_cats is None:
            pie_cats = chart_config["labels"]
        pie_series = chart_config.get("series")
        if pie_series:
            pie_values = pie_series[0]["data"]
        else:
            pie_values = chart_config["values"]
        cd.categories = [str(x) for x in pie_cats]
        cd.add_series("", pie_values)
    elif editorial:
        # Split the single series into N single-value series so each bar
        # carries its own colour (native charts have no per-point bar fill).
        data = chart_config["series"][0]["data"]
        cats = [str(x) for x in chart_config["x"]]
        cd.categories = cats
        for i, v in enumerate(data):
            row = [None] * len(data)
            row[i] = v
            cd.add_series(cats[i], row)
    else:
        cd.categories = [str(x) for x in chart_config["x"]]
        for s in chart_config["series"]:
            cd.add_series(s["name"], s["data"])

    chart_shape = slide.shapes.add_chart(ctype, left, top, width, height, cd)
    chart = chart_shape.chart

    if ctype == XL_CHART_TYPE.PIE:
        plot = chart.plots[0]
        for i, point in enumerate(plot.series[0].points):
            color = GREEN if i == accent_idx else NON_ACCENT_COLORS[
                min(i, len(NON_ACCENT_COLORS) - 1)
            ]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = color
    elif editorial:
        ramp = _green_ramp(len(list(chart.series)), accent_idx)
        chart.has_legend = False
        for i, s in enumerate(chart.series):
            fill = s.format.fill
            fill.solid()
            fill.fore_color.rgb = ramp[i]
            s.format.line.fill.background()
        try:
            chart.plots[0].gap_width = 60
            chart.plots[0].overlap = 100  # bars share the category slot
        except Exception:
            pass
    else:
        _apply_series_colors(chart, accent_idx)

    text_color = WHITE if dark else GRAPHITE

    if chart_config.get("title"):
        chart.has_title = True
        # Chart title lives inside the chart graphicFrame, which
        # apply_kpi_emphasis does not visit — ** here is a pure leak; strip it.
        chart.chart_title.text_frame.text = sanitize_text(chart_config["title"])
        _style_text(chart.chart_title.text_frame, color=text_color, size_pt=14, bold=True)
    else:
        chart.has_title = False

    if ctype != XL_CHART_TYPE.PIE and not editorial:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.legend.include_in_layout = False
        try:
            chart.legend.font.size = Pt(10)
            chart.legend.font.color.rgb = text_color
            chart.legend.font.name = "SB Sans Display"
        except Exception:
            pass

    if ctype == XL_CHART_TYPE.AREA_STACKED_100:
        try:
            chart.value_axis.tick_labels.number_format = "0%"
            chart.value_axis.tick_labels.font.size = Pt(10)
            chart.value_axis.tick_labels.font.color.rgb = text_color
            chart.category_axis.tick_labels.font.size = Pt(10)
            chart.category_axis.tick_labels.font.color.rgb = text_color
        except Exception:
            pass

    return chart_shape


def render_chart_pptx_slide(slide, chart_config, dark=False):
    """
    Полная сборка слайда с editable chart:
    title + native PowerPoint chart + caption (как chart_native, но редактируемая).

    chart_config:
      title: header слайда (str)
      slide_title: alias для title (для совместимости с chart_engine)
      caption: подпись под chart (str, optional)
      type / x / series / accent_idx / labels / values: см. add_chart_to_slide
    """
    from kpi_renderer import _add_text_box, set_slide_title
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    editorial = is_editorial_eligible(chart_config)
    if editorial:
        dark = True
    text_color = WHITE if dark else GRAPHITE

    # Title — в штатный placeholder шаблона, единый стиль (Problem #6)
    title_text = chart_config.get("title") or chart_config.get("slide_title", "")
    if title_text:
        set_slide_title(slide, title_text, dark=dark)

    if editorial:
        ZONE_X, ZONE_Y, ZONE_W, ZONE_H = 60, 150, 760, 450  # leave right gutter
    else:
        ZONE_X, ZONE_Y, ZONE_W, ZONE_H = 60, 120, 1160, 480
    chart_inner_cfg = {k: v for k, v in chart_config.items()
                       if k not in ("title", "slide_title", "caption")}

    add_chart_to_slide(
        slide, chart_inner_cfg,
        Emu(ZONE_X * 9525), Emu(ZONE_Y * 9525),
        Emu(ZONE_W * 9525), Emu(ZONE_H * 9525),
        dark=dark,
    )

    if editorial:
        # Large white выноска (peak value) to the right of the plot. Size the
        # font by length and disable wrap so multi-digit/decimal peaks like
        # "13.4" stay on one line instead of breaking (live deck_a.s5).
        data = chart_config["series"][0]["data"]
        acc = chart_config.get("accent_idx", -1)
        peak = data[acc] if 0 <= acc < len(data) else max(data)
        peak_s = str(peak)
        call_w = 1245 - (ZONE_X + ZONE_W + 20)  # use the full right gutter
        peak_font = 150 if len(peak_s) <= 3 else (115 if len(peak_s) <= 4 else 92)
        call_box = _add_text_box(slide, ZONE_X + ZONE_W + 20, 180, call_w, 360,
                      peak_s, font_size_pt=peak_font, bold=True, color=WHITE,
                      align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        call_box.text_frame.word_wrap = False

    caption = chart_config.get("caption", "")
    if caption:
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(0 * 9525), Emu((620 - 10) * 9525),
            Emu(1280 * 9525), Emu(80 * 9525),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)
        bg.line.fill.background()
        _add_text_box(slide, 35, 620, 1209, 60, caption,
                      font_size_pt=14, color=text_color,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


if __name__ == "__main__":
    import sys
    from pptx import Presentation

    out = sys.argv[1] if len(sys.argv) > 1 else "/tmp/chart_native_test.pptx"
    p = Presentation()
    blank = p.slide_layouts[6]
    slide = p.slides.add_slide(blank)
    cfg = {
        "type": "area_stacked",
        "title": "Структурные сдвиги занятости (1850–2050)",
        "x": [1850, 1875, 1900, 1925, 1950, 1975, 2000, 2025, 2050],
        "accent_idx": 3,
        "series": [
            {"name": "Сельское хозяйство", "data": [45, 40, 30, 20, 10, 5, 3, 2, 2]},
            {"name": "Промышленность", "data": [25, 30, 40, 45, 35, 25, 18, 12, 8]},
            {"name": "Услуги", "data": [25, 25, 25, 30, 50, 65, 70, 65, 50]},
            {"name": "Реляционный сектор", "data": [3, 3, 3, 3, 3, 3, 7, 18, 35]},
            {"name": "Другое", "data": [2, 2, 2, 2, 2, 2, 2, 3, 5]},
        ],
    }
    add_chart_to_slide(slide, cfg, Emu(457200), Emu(914400), Emu(8229600), Emu(4572000))
    p.save(out)
    print(f"Saved {out}")
