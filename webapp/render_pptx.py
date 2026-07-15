"""Pack a saved deck into a .pptx where each slide is one full-bleed 16:9 image.

We already render every slide to a 1920x1080 PNG for the PNG-ZIP export; a PPTX
that drops one such image per slide onto a blank 16:9 layout reproduces the deck
1-to-1 (fonts, brand, charts — exactly as the HTML renders). The trade-off is that
text inside a slide is a picture, not editable PowerPoint text — acceptable for a
"looks identical" hand-off. Reuses render_png.render_slides for the heavy step.
"""
from __future__ import annotations

from pathlib import Path

from webapp import render_png

# 16:9 at PowerPoint's EMU scale (914400 EMU = 1 inch). 13.333in x 7.5in is the
# canonical widescreen slide; the 1920x1080 PNG shares that aspect, so a full-bleed
# placement neither letterboxes nor distorts.
_EMU_PER_INCH = 914400
_SLIDE_W = int(13.333 * _EMU_PER_INCH)
_SLIDE_H = int(7.5 * _EMU_PER_INCH)


def build_pptx(pngs: dict[int, Path], out_pptx: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    prs.slide_width = Emu(_SLIDE_W)
    prs.slide_height = Emu(_SLIDE_H)
    blank = prs.slide_layouts[6]  # the built-in "Blank" layout — no placeholders
    for index in sorted(pngs):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(pngs[index]), 0, 0,
                                 width=Emu(_SLIDE_W), height=Emu(_SLIDE_H))
    prs.save(str(out_pptx))
    return out_pptx


def export_pptx(deck_html: Path, out_pptx: Path,
                pngs: dict[int, Path] | None = None) -> Path:
    if pngs is None:
        pngs = render_png.render_slides(deck_html)
    return build_pptx(pngs, out_pptx)
