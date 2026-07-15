import io

from PIL import Image
from pptx import Presentation

import webapp.render_pptx as render_pptx


def _png(path, size=(1920, 1080)):
    Image.new("RGB", size, (10, 20, 30)).save(path)


def test_build_pptx_one_full_bleed_image_per_slide(tmp_path):
    a = tmp_path / "s1.png"
    b = tmp_path / "s2.png"
    _png(a)
    _png(b)
    out = tmp_path / "deck.pptx"
    render_pptx.build_pptx({1: a, 2: b}, out)

    prs = Presentation(str(out))
    assert len(prs.slides) == 2
    # 16:9 widescreen canvas
    assert prs.slide_width == render_pptx._SLIDE_W
    assert prs.slide_height == render_pptx._SLIDE_H
    for slide in prs.slides:
        pics = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
        assert len(pics) == 1
        pic = pics[0]
        # full-bleed: top-left origin, fills the whole slide
        assert pic.left == 0 and pic.top == 0
        assert pic.width == render_pptx._SLIDE_W
        assert pic.height == render_pptx._SLIDE_H


def test_export_pptx_reuses_supplied_pngs_without_rendering(monkeypatch, tmp_path):
    deck = tmp_path / "deck.html"
    deck.write_text('<section class="slide">a</section>', encoding="utf-8")
    png = tmp_path / "s1.png"
    _png(png)

    def boom(*a, **k):
        raise AssertionError("render_slides must not run when pngs are supplied")

    monkeypatch.setattr(render_pptx.render_png, "render_slides", boom)
    out = render_pptx.export_pptx(deck, tmp_path / "out.pptx", pngs={1: png})
    assert out.is_file()
    assert len(Presentation(str(out)).slides) == 1


def test_export_pptx_renders_when_pngs_missing(monkeypatch, tmp_path):
    deck = tmp_path / "deck.html"
    deck.write_text('<section class="slide">a</section>'
                    '<section class="slide">b</section>', encoding="utf-8")

    def fake_render(html_path):
        paths = {}
        for i in (1, 2):
            p = tmp_path / f"qa-{i}.png"
            _png(p)
            paths[i] = p
        return paths

    monkeypatch.setattr(render_pptx.render_png, "render_slides", fake_render)
    out = render_pptx.export_pptx(deck, tmp_path / "out.pptx")
    assert len(Presentation(str(out)).slides) == 2
