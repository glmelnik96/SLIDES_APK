"""Render-free обогащение parse_pptx: текст из групп, числа из диаграмм, заметки.

python-pptx читает pptx прямо из XML — рендер (LibreOffice) не нужен. Раньше
parse_pptx терял три вещи: (1) текст внутри сгруппированных фигур, (2) данные
диаграмм, (3) заметки докладчика. Тесты строят НАСТОЯЩИЙ pptx (python-pptx умеет
и создавать группы/графики) и проверяют, что контент доезжает до InputDoc.

Ключевой инвариант: заметки идут в Section.notes, а НЕ в blocks — иначе они
утекут прямо на слайд в режиме «Точный перенос» (exact рендерит heading+blocks).
"""
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from htmlslides.parsers.pptx import parse_pptx


def _fixture(tmp_path: Path) -> Path:
    """3 слайда: [1] группа с текстом, [2] диаграмма, [3] видимый текст + заметка."""
    prs = Presentation()
    blank = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank)
    grp = s1.shapes.add_group_shape()
    tb = grp.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.text = "Текст-внутри-группы"

    s2 = prs.slides.add_slide(blank)
    data = CategoryChartData()
    data.categories = ["2023", "2024", "2025"]
    data.add_series("Выручка", (10, 20, 35))
    s2.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(1), Inches(1), Inches(6), Inches(4), data)

    s3 = prs.slides.add_slide(blank)
    box = s3.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Видимый текст слайда"
    s3.notes_slide.notes_text_frame.text = "СЕКРЕТНАЯ-ЗАМЕТКА"

    out = tmp_path / "fixture.pptx"
    prs.save(str(out))
    return out


def _all_text(section) -> str:
    """Весь текст из блоков секции (для проверки «контент дошёл»)."""
    parts: list[str] = []
    for b in section.blocks:
        parts.append(getattr(b, "text", "") or "")
        parts.extend(getattr(b, "items", []) or [])
        for row in getattr(b, "rows", []) or []:
            parts.extend(row)
    return " ".join(parts)


def test_group_text_is_captured(tmp_path):
    """Текст внутри сгруппированных фигур раньше терялся — теперь собирается."""
    doc = parse_pptx(_fixture(tmp_path))
    assert "Текст-внутри-группы" in _all_text(doc.sections[0])


def test_chart_numbers_captured_as_table(tmp_path):
    """Диаграмма отдаётся таблицей «категория × ряд»: цифры доходят до планировщика."""
    text = _all_text(parse_pptx(_fixture(tmp_path)).sections[1])
    assert "Выручка" in text                    # имя ряда
    assert "2025" in text and "35" in text       # категория и её значение


def test_notes_go_to_section_notes_not_blocks(tmp_path):
    """Заметка — в Section.notes и НЕ в контенте (иначе утечёт в exact)."""
    s3 = parse_pptx(_fixture(tmp_path)).sections[2]
    assert s3.notes == "СЕКРЕТНАЯ-ЗАМЕТКА"
    assert "СЕКРЕТНАЯ-ЗАМЕТКА" not in _all_text(s3)


def test_exact_mode_does_not_leak_notes(tmp_path):
    """Контракт exact: рендерятся heading+blocks, заметки докладчика — никогда."""
    from htmlslides.pipeline.exact_builder import build_exact_plan
    doc = parse_pptx(_fixture(tmp_path))
    plan, _ = build_exact_plan(doc)
    html = " ".join(s.content.get("html", "") for s in plan.slides)
    assert "Видимый текст слайда" in html        # контент на месте
    assert "СЕКРЕТНАЯ-ЗАМЕТКА" not in html        # заметка не утекла


def test_notes_reach_planner_section_text(tmp_path):
    """Заметки доезжают в текст для планировщика — как контекст для дизайна."""
    from htmlslides.pipeline.planner import _section_to_text
    txt = _section_to_text(parse_pptx(_fixture(tmp_path)).sections[2])
    assert "СЕКРЕТНАЯ-ЗАМЕТКА" in txt


# 1x1 прозрачный PNG — валидный blob для add_picture
_PNG = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
    b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\nIDATx\x9cc\x00\x01"
    b"\x00\x00\x05\x00\x01\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
)


def _fixture_alien_picture(tmp_path: Path) -> Path:
    """pptx, где у картинки content-type НЕ из реестра python-pptx.

    Репро прод-файла 2026-07-29 (job 22/23): дека со сторонним экспортом несла
    85 картинок с content-type `image/svg` (не стандартный `image/svg+xml`).
    python-pptx собирает такую часть как generic Part, у которого нет `.image`,
    и `shape.image` падает AttributeError. Строим то же самое: настоящий pptx с
    картинкой + текстом, затем в [Content_Types].xml переписываем тип png-части.
    """
    import io
    import zipfile

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(_PNG), Inches(1), Inches(1))
    box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1))
    box.text_frame.text = "Текст-рядом-с-картинкой"
    normal = tmp_path / "normal.pptx"
    prs.save(str(normal))

    out = tmp_path / "alien.pptx"
    with zipfile.ZipFile(normal) as src, \
            zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            data = src.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(b'ContentType="image/png"',
                                    b'ContentType="image/svg"')
            dst.writestr(item, data)
    return out


def test_alien_picture_mime_does_not_crash_parse(tmp_path):
    """Картинка с типом вне реестра python-pptx не роняет parse_pptx.

    Прод 2026-07-29: файл падал за 0,5 с с «Внутренняя ошибка сборки» — до
    единого LLM-вызова. Контракт тот же, что у _chart_blocks: парсер не падает,
    текст слайда доезжает, нечитаемая картинка в худшем случае пропускается."""
    doc = parse_pptx(_fixture_alien_picture(tmp_path))
    assert "Текст-рядом-с-картинкой" in _all_text(doc.sections[0])


def _fixture_soft_break(tmp_path: Path) -> Path:
    """Слайд, где в заголовке стоит МЯГКИЙ перенос (<a:br/>), а не новый абзац.

    В pptx это обычный приём вёрстки: автор ломает длинную фразу по строкам
    внутри одного абзаца. python-pptx отдаёт такой перенос как \\x0b прямо
    внутри `shape.text`."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])   # Title Only
    slide.shapes.title.text_frame.text = "Миграция,\x0bрезервное копирование"
    box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1))
    box.text_frame.text = "Этап 1:\x0bперенос данных"
    out = tmp_path / "softbreak.pptx"
    prs.save(str(out))
    return out


def test_soft_break_becomes_a_space(tmp_path):
    """Мягкий перенос — это пробел, а не символ контента.

    Раньше \\x0b ехал в модель как есть: заголовок деки резался посреди фразы
    («Миграция,» + невидимый управляющий символ), а на обложке хвост уходил в
    подзаголовок отдельным предложением. Блок и заголовок у нас — ОДНА строка."""
    doc = parse_pptx(_fixture_soft_break(tmp_path))
    section = doc.sections[0]
    assert section.heading == "Миграция, резервное копирование"
    assert doc.title == "Миграция, резервное копирование"
    text = _all_text(section)
    assert "Этап 1: перенос данных" in text
    assert "\x0b" not in section.heading and "\x0b" not in text


def test_alien_picture_blob_still_captured(tmp_path):
    """Blob сторонней картинки достаём напрямую из part — контент не теряется,
    а нестандартный `image/svg` нормализуется в честный `image/svg+xml`."""
    from htmlslides.parsers.base import ImageBlock
    doc = parse_pptx(_fixture_alien_picture(tmp_path))
    images = [b for b in doc.sections[0].blocks if isinstance(b, ImageBlock)]
    assert images and images[0].data == _PNG
    assert images[0].mime == "image/svg+xml"
