"""«Стеклянная сборка»: скоринг кандидатов, пороги сомнения, степпер, ответы.

Fake-клиент вместо LLM (паттерн test_diagram_filler.py) и fake fill_slide
вместо заполнения: проверяем обвязку — фильтрацию кандидатов, фолбэк как
«место сомнения» (confidence 0.0), пропуск needs_input степпером, доводку
слайда по ответу и API-контракты /glass/*.

Контракт после переработки (2026-08-20, жалобы «сайт лежит на старте» и
«вопросы не видны»): старт — parse-only без единого вызова модели (слайды
уходят в status="unscored"), разметка ленивая — score_next/step_fill по одному
слайду, вопрос НЕ блокирует конвейер (слайд откладывается, сборка идёт дальше),
каждый шаг возвращает action: "score"|"fill"|None — клиент останавливает цикл
на None, а done значит строго «всё заполнено и вопросов нет».
"""
import os
os.environ["SLIDES_APP_SKIP_SHIM"] = "1"

import re

import pytest
from fastapi.testclient import TestClient

import webapp.app as appmod
import webapp.config as cfg
from htmlslides.library import SlotValidationError, TemplateLibrary
from htmlslides.parsers.base import ListBlock, Section, TextBlock
from webapp import draft, glass
from webapp.glass import Candidate, SectionChoice


def _client_app(monkeypatch, tmp_path):
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    monkeypatch.setattr(cfg.settings, "db_url",
                        f"sqlite+aiosqlite:///{tmp_path / 't.db'}")
    monkeypatch.setattr(cfg.settings, "data_dir", tmp_path)
    monkeypatch.setattr(cfg.settings, "dev_user_id", "")
    return TestClient(appmod.app)


def H(uid="u1"):
    return {"X-User-Id": uid}


@pytest.fixture(scope="module")
def library():
    return TemplateLibrary.load()


class FakeClient:
    """Отдаёт заготовленные ответы по очереди; Exception в очереди — бросается."""

    def __init__(self, replies):
        self.replies = list(replies)
        self.calls = []

    def chat_json(self, messages, model_cls, *, max_tokens=4096,
                  retries=2, extra_body=None):
        self.calls.append({"messages": messages, "extra_body": extra_body})
        reply = self.replies.pop(0)
        if isinstance(reply, Exception):
            raise reply
        return reply


def _section(heading="Процесс", text="Заявка проходит согласование и запуск."):
    return Section(heading=heading, level=2, blocks=[TextBlock(text=text)])


def _fake_fill(client, library, sp, *, deck_title="", extra="", diagram_kind=""):
    return sp.model_copy(update={
        "content": {"title": f"Слайд {sp.index}", "brief": sp.content.get("brief", "")}})


# ── plan_section_candidates ──────────────────────────────────────────────────
def test_candidates_valid_reply_filtered_to_known(library):
    fake = FakeClient([SectionChoice(
        candidates=[Candidate(template_id="three-col", confidence=0.9),
                    Candidate(template_id="bogus", confidence=0.8),
                    Candidate(template_id="timeline", confidence=0.5)],
        question="  Что важнее?  ")])
    out = glass.plan_section_candidates(fake, library, _section())
    assert [c.template_id for c in out.candidates] == ["three-col", "timeline"]
    assert out.question == "Что важнее?"
    # скоринг — текст-в-JSON без reasoning, как слот-филл
    assert fake.calls[0]["extra_body"] == {"thinking": {"type": "disabled"}}


def test_candidates_failure_falls_back_with_zero_confidence(library):
    """Сбой сети/формата → эвристика planner с confidence 0.0 — честное
    «место сомнения», а не тихий уверенный выбор."""
    for fake in (FakeClient([RuntimeError("boom")]),
                 FakeClient([SectionChoice(          # все id чужие
                     candidates=[Candidate(template_id="ghost", confidence=1.0)])])):
        out = glass.plan_section_candidates(fake, library, _section())
        assert len(out.candidates) == 1
        assert out.candidates[0].confidence == 0.0
        assert out.candidates[0].template_id in {t.id for t in library.templates}


def test_candidates_menu_excludes_system_templates(library):
    """Обложки/контакты/дивайдеры ставит система — в меню раздела их нет."""
    menu = glass._menu(library)
    for tid in ("cover", "contacts", "back-cover", "blank"):
        assert f"- {tid} " not in menu
    assert "- three-col " in menu and "- diagram " in menu


def test_candidates_collapse_cosmetic_twins(library):
    """statement и statement-green — один макет в двух цветах. Двумя чипами они
    предлагают выбор без выбора, а их близкая уверенность поднимала ложный
    вопрос: разрыв меньше GAP_FLOOR считался сомнением, хотя различался только
    фон. Схлопываем в один вариант — и вопрос задаётся по существу."""
    fake = FakeClient([SectionChoice(
        candidates=[Candidate(template_id="statement", confidence=0.86),
                    Candidate(template_id="statement-green", confidence=0.8),
                    Candidate(template_id="quote", confidence=0.4)])])
    out = glass.plan_section_candidates(fake, library, _section())
    assert [c.template_id for c in out.candidates] == ["statement", "quote"]
    assert glass._doubtful(out) is False   # без дубля разрыв 0.46 — вопроса нет


def test_doubtful_thresholds():
    ok = lambda *pairs: SectionChoice(candidates=[
        Candidate(template_id="x", confidence=c) for c in pairs])
    assert glass._doubtful(ok(0.9)) is False
    assert glass._doubtful(ok(0.5)) is True              # ниже пола 0.6
    assert glass._doubtful(ok(0.8, 0.7)) is True         # разрыв < 0.15
    assert glass._doubtful(ok(0.9, 0.5)) is False
    assert glass._doubtful(SectionChoice()) is True      # пусто = сомнение


# ── start_glass ──────────────────────────────────────────────────────────────
def test_start_glass_is_parse_only_no_llm(monkeypatch, tmp_path):
    """Старт мгновенный: разбор документа БЕЗ единого вызова модели.

    Жалоба 2026-08-20: /glass/start скорил ВСЕ разделы в одном HTTP-запросе —
    «сайт по факту лежит» 5-20 минут до первого ответа. Теперь старт только
    режет документ на слайды-заготовки (status="unscored"), временный макет
    ставит эвристика планировщика, а разметка уходит в ленивые score_next
    по ходу конвейера."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    src = tmp_path / "doc.md"
    src.write_text("# Отчёт\n\n## Уверенный раздел\n\nТекст про метрики.\n\n"
                   "## Спорный раздел\n\nМало контента.\n", encoding="utf-8")
    monkeypatch.setattr(glass, "_kimi", lambda: (_ for _ in ()).throw(
        AssertionError("старт не имеет права звать модель")))
    plan = glass.start_glass("s1", src)

    assert plan.title == "Отчёт"
    cover = plan.slides[0]
    assert cover.template_id == "cover" and cover.filled
    assert len(plan.slides) == 3
    for s in plan.slides[1:]:
        assert s.status == "unscored" and not s.filled
        assert s.brief
        assert s.candidates is None and s.question is None
        assert s.template_id                       # временный макет эвристики
    # DeckPlan-as-truth: план сохранён, дека отрисована из него
    assert draft.load_plan("s1").model_dump() == plan.model_dump()
    from webapp.paths import session_dir
    assert (session_dir("s1") / "deck.html").is_file()


def test_cover_keeps_the_whole_document_title(monkeypatch, tmp_path):
    """Длинное название не режется посреди слова: заголовок обложки узкий, хвост
    уходит в подзаголовок (иначе первый слайд деки — обрубок и текст-рыба)."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    src = tmp_path / "doc.md"
    long_title = "Программа модернизации логистического оператора «Северный путь»"
    src.write_text(f"# {long_title}\n\n## Раздел\n\nТекст про метрики.\n",
                   encoding="utf-8")
    plan = glass.start_glass("cov", src)

    from htmlslides.library import TemplateLibrary
    cap = TemplateLibrary.load().get("cover").slots["title"].max_chars
    cover = plan.slides[0].content
    assert len(cover["title"]) <= cap
    assert long_title[len(cover["title"])] == " "   # разрез по границе слова
    assert cover["title"] + " " + cover["subtitle"] == long_title


def test_cover_short_title_ships_without_sample_subtitle(monkeypatch, tmp_path):
    """Готовый слайд не досочиняется примерами из каталога: у короткого названия
    подзаголовка нет, и раньше обложка уезжала автору с «Коротким подзаголовком»
    под ним — текст-рыба на первом же слайде деки."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    src = tmp_path / "doc.md"
    src.write_text("# Итоги\n\n## Раздел\n\nТекст про метрики.\n", encoding="utf-8")
    glass.start_glass("sub", src)
    from webapp.paths import session_dir
    html = (session_dir("sub") / "deck.html").read_text(encoding="utf-8")
    # Смотрим ТОЛЬКО обложку: у ещё не заполненных слайдов рыба в пустых слотах
    # законна — это превью макета, пока ИИ до него не дошёл.
    cover = re.split(r"<section[^>]*\bclass=\"slide", html)[1]
    assert "Итоги" in cover
    assert "Короткий подзаголовок" not in cover


def test_deck_title_falls_back_to_the_uploaded_file_name(monkeypatch, tmp_path):
    """У pptx без титульного текста заголовок берётся из имени файла АВТОРА:
    на диске загрузка всегда `input.<ext>`, и обложка выходила с «INPUT»."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    src = tmp_path / "input.md"
    src.write_text("## Раздел\n\nТекст про метрики.\n", encoding="utf-8")
    plan = glass.start_glass("t1", src,
                             origin_name="Партнёры_Инфраструктура_для_ИИ.pptx")
    assert plan.title == "Партнёры Инфраструктура для ИИ"   # «_» — это пробелы

    # имени файла нет вовсе → заголовок первого раздела, но никогда не «input»
    plan = glass.start_glass("t2", src)
    assert plan.title == "Раздел"


def test_image_only_section_asks_instead_of_inventing(monkeypatch, tmp_path):
    """Раздел из одних картинок уходит в needs_input при скоринге, даже если
    ИИ уверен в макете.

    Стеклянная сборка идёт мимо vision-ветки чёрного ящика: текста у такого
    раздела нет, и филлер не «собирал хуже» — он ДОСОЧИНЯЛ тезисы и цифры,
    которых в документе нет. Спросить автора честнее, чем выдумать."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    src = tmp_path / "doc.md"
    src.write_text("# Дека\n\n## Архитектура\n\n![](a.png)\n\n![](b.png)\n",
                   encoding="utf-8")
    glass.start_glass("img", src)
    out = glass.score_next("img", client=FakeClient([
        SectionChoice(candidates=[Candidate(template_id="stats-row",
                                            confidence=0.95),
                                  Candidate(template_id="three-col",
                                            confidence=0.2)])]))
    assert out["action"] == "score" and out["index"] == 2
    s = draft.load_plan("img").slides[1]
    assert s.status == "needs_input"
    assert "только изображения" in (s.question or "")
    assert s.candidates == ["stats-row", "three-col"]   # чипы на месте


_PNG = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
    b"\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\nIDATx\x9cc\x00\x01"
    b"\x00\x00\x05\x00\x01\r\n-\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
)


def _pptx_with_blind_slide(tmp_path):
    """pptx из двух слайдов: [1] текстовый, [2] одна картинка без текста."""
    from pptx import Presentation
    from pptx.util import Inches
    import io

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    s1.shapes.title.text_frame.text = "Как устроена платформа"
    box = s1.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(1))
    box.text_frame.text = ("Платформа собрана из трёх слоёв: хранение, "
                           "вычисления и оркестрация задач.")
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    s2.shapes.add_picture(io.BytesIO(_PNG), Inches(1), Inches(1),
                          Inches(6), Inches(4))
    out = tmp_path / "input.pptx"
    prs.save(str(out))
    return out


def _blind_session(monkeypatch, tmp_path, sid):
    """pptx со слепым слайдом кладётся в сессию как input.pptx — так делает
    /glass/start, и так ленивый vision находит исходник при скоринге."""
    from webapp.paths import session_dir
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    dest = session_dir(sid) / "input.pptx"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(_pptx_with_blind_slide(tmp_path).read_bytes())
    return dest


def test_blind_pptx_section_is_read_from_the_source_slide(monkeypatch, tmp_path):
    """Раздел без текста дочитывается с картинки исходного слайда — ЛЕНИВО,
    при скоринге именно этого раздела, а не на старте за все сразу.

    Спросить автора честнее, чем выдумать, — но ещё честнее прочитать: у pptx
    исходный слайд можно отрендерить и показать модели. Прочитанное едет в бриф
    С ПОМЕТКОЙ (автор видит в панели контекста, откуда текст), раздел перестаёт
    быть слепым, и вопрос на нём не поднимается."""
    dest = _blind_session(monkeypatch, tmp_path, "blind")
    png = tmp_path / "slide-02.png"
    png.write_bytes(_PNG)
    shots = tmp_path / "shot-01.png"
    shots.write_bytes(_PNG)
    monkeypatch.setattr("htmlslides.parsers.render.render_pptx_pngs",
                        lambda src, out, **kw: [shots, png])

    glass.start_glass("blind", dest)
    read = glass._SlideRead(text="Схема потоков данных: источники, шина Kafka, "
                                 "витрины. Внизу подпись «до 40 ТБ в сутки».")
    ok = SectionChoice(candidates=[Candidate(template_id="diagram",
                                             confidence=0.9),
                                   Candidate(template_id="three-col",
                                             confidence=0.4)])
    text_client = FakeClient([ok])
    glass.score_next("blind", client=text_client)       # текстовый раздел
    blind_client = FakeClient([read, ok])
    glass.score_next("blind", client=blind_client)      # слепой: сперва читаем

    blind = draft.load_plan("blind").slides[2]
    assert blind.status is None                     # больше не слепой
    assert "Kafka" in blind.brief
    assert glass._IMAGE_READ_NOTE in blind.brief    # источник текста назван
    # читаем ТОЛЬКО слепой раздел: текстовый слайд рендерить и смотреть незачем
    assert sum(1 for c in text_client.calls
               if isinstance(c["messages"][-1]["content"], list)) == 0
    assert sum(1 for c in blind_client.calls
               if isinstance(c["messages"][-1]["content"], list)) == 1


def test_unreadable_source_still_asks_the_author(monkeypatch, tmp_path):
    """Нет LibreOffice (или картинка не прочиталась) — возвращаемся к вопросу,
    а не к выдумке: гард слепого раздела остаётся последней линией."""
    dest = _blind_session(monkeypatch, tmp_path, "blind2")
    from htmlslides.parsers.render import RenderUnavailable

    def boom(src, out, **kw):
        raise RenderUnavailable("LibreOffice (soffice) not found")

    monkeypatch.setattr("htmlslides.parsers.render.render_pptx_pngs", boom)
    ok = SectionChoice(candidates=[Candidate(template_id="three-col",
                                             confidence=0.9)])
    glass.start_glass("blind2", dest)
    glass.score_next("blind2", client=FakeClient([ok]))
    glass.score_next("blind2", client=FakeClient([ok]))
    s = draft.load_plan("blind2").slides[2]
    assert s.status == "needs_input"
    assert "только изображения" in (s.question or "")


def test_image_section_with_text_goes_through(monkeypatch, tmp_path):
    """Картинка РЯДОМ с текстом — обычный раздел: гард не должен ловить всё,
    где вообще есть изображение, иначе вопросы посыплются на ровном месте."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    src = tmp_path / "doc.md"
    src.write_text("# Дека\n\n## Архитектура\n\n![](a.png)\n\nПлатформа собрана "
                   "из трёх слоёв: хранение, вычисления и оркестрация задач.\n",
                   encoding="utf-8")
    glass.start_glass("img2", src)
    glass.score_next("img2", client=FakeClient([
        SectionChoice(candidates=[Candidate(template_id="stats-row",
                                            confidence=0.95),
                                  Candidate(template_id="three-col",
                                            confidence=0.2)])]))
    assert draft.load_plan("img2").slides[1].status is None


def test_same_layout_never_runs_three_in_a_row(monkeypatch, tmp_path):
    """Однородный документ не превращается в деку из одного макета.

    На реальном прогоне СЕМЬ разделов подряд получили cards-6 — 45% деки одной
    сеткой. Разводим повторы вторым кандидатом самой модели, но только когда он
    почти так же хорош. Скоринг теперь ленивый и последовательный — каждый
    score_next видит уже выбранные макеты соседей."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    src = tmp_path / "doc.md"
    src.write_text("# Дека\n\n" + "".join(
        f"## Раздел {i}\n\nТри направления работы: люди, процессы и данные.\n\n"
        for i in range(1, 8)), encoding="utf-8")
    fake = FakeClient([
        SectionChoice(candidates=[Candidate(template_id="cards-6", confidence=0.9),
                                  Candidate(template_id="three-col", confidence=0.8)])
        for _ in range(7)])
    glass.start_glass("mono", src)
    for _ in range(7):
        glass.score_next("mono", client=fake)
    used = [s.template_id for s in draft.load_plan("mono").slides[1:]]
    assert len(used) == 7
    assert "three-col" in used
    run = max_run = 1
    for a, b in zip(used, used[1:]):
        run = run + 1 if a == b else 1
        max_run = max(max_run, run)
    assert max_run < 3, used


def test_variety_never_costs_a_much_worse_layout(monkeypatch, tmp_path):
    """Разнообразие не стоит плохого макета: если второй кандидат заметно хуже,
    повтор остаётся — лучше ровная дека, чем разная и мимо смысла."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    src = tmp_path / "doc.md"
    src.write_text("# Дека\n\n" + "".join(
        f"## Раздел {i}\n\nТри направления работы: люди, процессы и данные.\n\n"
        for i in range(1, 5)), encoding="utf-8")
    fake = FakeClient([
        SectionChoice(candidates=[Candidate(template_id="cards-6", confidence=0.95),
                                  Candidate(template_id="quote", confidence=0.1)])
        for _ in range(4)])
    glass.start_glass("mono2", src)
    for _ in range(4):
        glass.score_next("mono2", client=fake)
    used = [s.template_id for s in draft.load_plan("mono2").slides[1:]]
    assert used == ["cards-6"] * 4


def test_score_names_the_scoring_failure(monkeypatch, tmp_path):
    """Сбой скоринга не маскируется под вопрос ИИ: у слайда честный текст про
    осечку модели, а не «раздел ложится в несколько макетов» при одном чипе."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    src = tmp_path / "doc.md"
    src.write_text("# Т\n\n## Раздел\n\nТекст.\n", encoding="utf-8")
    glass.start_glass("s2", src)
    out = glass.score_next("s2", client=FakeClient([RuntimeError("boom")]))
    assert out["action"] == "score" and out["index"] == 2
    s = draft.load_plan("s2").slides[1]
    assert s.status == "needs_input"
    assert s.question == glass._FALLBACK_QUESTION
    assert len(s.candidates) == 1


def test_score_next_without_unscored_slides_is_a_noop(monkeypatch, tmp_path):
    """Скоринг на размеченном плане отвечает action=None — клиентский разведчик
    останавливает свой цикл, а не молотит модель по кругу."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    draft.save_plan("s2b", _outline_plan(question=False))
    out = glass.score_next("s2b", client=FakeClient([]))
    assert out["action"] is None and out["index"] is None
    assert out["unscored"] == [] and out["done"] is False


def test_question_text_matches_the_number_of_options():
    """Заготовка вопроса зависит от числа вариантов: «несколько макетов» при
    одном чипе — ложь, из-за которой выбор выглядел сломанным."""
    two = SectionChoice(candidates=[Candidate(template_id="a", confidence=0.5),
                                    Candidate(template_id="b", confidence=0.45)])
    one = SectionChoice(candidates=[Candidate(template_id="a", confidence=0.4)])
    assert glass._question_for(two) == glass._DEFAULT_QUESTION
    assert glass._question_for(one) == glass._ONE_CANDIDATE_QUESTION
    one.question = "  Что показать?  "        # свой вопрос модели важнее заготовки
    assert glass._question_for(one) == "Что показать?"


# ── step_fill ────────────────────────────────────────────────────────────────
def _outline_plan(*, question=True, at=2):
    """Размеченный аутлайн (все слайды прошли скоринг): готов к заполнению.

    question=False — без вопроса; at — номер слайда с вопросом. Вопрос сборку
    не останавливает — спорный слайд просто откладывается, а вопрос в хвосте
    оставляет шагам работу перед собой (нужно для гонок)."""
    plan = draft.DraftPlan(title="Т", slides=[
        draft.DraftSlide(template_id="cover", content={"title": "Т"}, filled=True),
        draft.DraftSlide(template_id="three-col", brief="спорный"),
        draft.DraftSlide(template_id="stats-row", brief="метрики"),
        draft.DraftSlide(template_id="timeline", brief="этапы"),
    ])
    if question:
        plan.slides[at - 1] = plan.slides[at - 1].model_copy(update={
            "status": "needs_input", "question": "Какой макет?",
            "candidates": ["three-col", "timeline"]})
    return plan


def test_question_does_not_block_the_conveyor(monkeypatch, tmp_path):
    """Вопрос НЕ останавливает сборку: спорный слайд откладывается, остальная
    лента заполняется, ответить можно в любой момент.

    Жалоба 2026-08-20: степпер вставал на первом needs_input и молча ждал —
    автор не видел вопроса, «сборка зависла», хвост деки так и оставался
    рыбой. Теперь конвейер обходит вопрос, а done наступает только когда
    ни работы, ни открытых вопросов не осталось."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler
    monkeypatch.setattr(filler, "fill_slide", _fake_fill)
    draft.save_plan("s3", _outline_plan())        # вопрос на слайде 2

    out1 = glass.step_fill("s3", client=FakeClient([]))
    assert out1["action"] == "fill" and out1["index"] == 3
    out2 = glass.step_fill("s3", client=FakeClient([]))
    assert out2["action"] == "fill" and out2["index"] == 4
    out3 = glass.step_fill("s3", client=FakeClient([]))
    assert out3["action"] is None and out3["index"] is None
    assert out3["done"] is False and out3["open_questions"] == [2]
    assert draft.load_plan("s3").slides[1].filled is False   # ждёт автора

    glass.answer("s3", 2, template_id="timeline")
    out4 = glass.step_fill("s3", client=FakeClient([]))
    assert out4["action"] == "fill" and out4["index"] == 2
    out5 = glass.step_fill("s3", client=FakeClient([]))
    assert out5["action"] is None and out5["done"] is True
    plan = draft.load_plan("s3")
    assert all(s.filled for s in plan.slides)


def test_step_scores_when_nothing_is_ready_to_fill(monkeypatch, tmp_path):
    """Шаг без готовых слайдов не простаивает: сам размечает следующий
    unscored — конвейер жив, даже если клиентский разведчик отстал."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler
    monkeypatch.setattr(filler, "fill_slide", _fake_fill)
    draft.save_plan("sc1", draft.DraftPlan(title="Т", slides=[
        draft.DraftSlide(template_id="cover", content={"title": "Т"}, filled=True),
        draft.DraftSlide(template_id="three-col", brief="тема", status="unscored"),
    ]))

    out1 = glass.step_fill("sc1", client=FakeClient([
        SectionChoice(candidates=[Candidate(template_id="stats-row",
                                            confidence=0.9)])]))
    assert out1["action"] == "score" and out1["index"] == 2
    assert out1["unscored"] == [] and out1["done"] is False
    s = draft.load_plan("sc1").slides[1]
    assert s.status is None and s.template_id == "stats-row"

    out2 = glass.step_fill("sc1", client=FakeClient([]))
    assert out2["action"] == "fill" and out2["index"] == 2
    assert out2["done"] is True


def test_step_prefers_filling_over_scoring(monkeypatch, tmp_path):
    """Готовый слайд важнее разметки следующего: автор видит контент раньше,
    а разведку добирает параллельный /glass/score."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler
    monkeypatch.setattr(filler, "fill_slide", _fake_fill)
    plan = _outline_plan(question=False)
    plan.slides.append(draft.DraftSlide(template_id="three-col", brief="хвост",
                                        status="unscored"))
    draft.save_plan("sc2", plan)

    out = glass.step_fill("sc2", client=FakeClient([]))
    assert out["action"] == "fill" and out["index"] == 2
    assert out["unscored"] == [5]


def test_step_fill_survives_fill_failure(monkeypatch, tmp_path):
    """Осечка модели не роняет шаг, но и не выдаёт себя за успех.

    Пустой content нельзя оставлять: draft_render дорисовывает пустые слоты
    образцами из каталога (для diagram — целой блок-схемой), и автор принял бы
    выдумку за свой текст. Деградируем на blank с темой в заголовке — ровно как
    чёрный ящик — и помечаем слайд failed, чтобы шаг был виден в панели."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler

    def _boom(client, library, sp, **kw):
        raise filler.FillError("slide 2: сеть упала")
    monkeypatch.setattr(filler, "fill_slide", _boom)
    draft.save_plan("s4", _outline_plan(question=False))

    out = glass.step_fill("s4", client=FakeClient([]))
    assert out["action"] == "fill" and out["index"] == 2
    assert out["failed"] == [2]
    s = draft.load_plan("s4").slides[1]
    assert s.filled and s.status == "failed"
    assert s.template_id == "blank" and s.content.get("title")
    assert s.brief                       # тема цела — слайд можно дописать руками
    # Подмену макета объясняем словами: молчащая заглушка читалась как
    # «выбрал один макет, применился другой».
    assert s.question and "не удалось заполнить макет" in s.question.lower()


def test_failed_slide_does_not_pause_the_build(monkeypatch, tmp_path):
    """Осечка — не вопрос: её чинят в редакторе, а сборка идёт дальше.

    Этап 3 переработки 2026-08-20: перед заглушкой у заполнения есть один
    ретрай, поэтому «настоящая» осечка — это ДВА сбоя подряд на одном слайде."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler
    calls = []

    def _boom_twice(client, library, sp, **kw):
        calls.append(sp.index)
        if len(calls) <= 2:
            raise filler.FillError("сеть упала")
        return _fake_fill(client, library, sp, **kw)
    monkeypatch.setattr(filler, "fill_slide", _boom_twice)
    draft.save_plan("s4b", _outline_plan(question=False))

    out1 = glass.step_fill("s4b", client=FakeClient([]))
    assert calls == [2, 2]               # осечка + ретрай на том же слайде
    assert out1["failed"] == [2] and out1["action"] == "fill"
    out2 = glass.step_fill("s4b", client=FakeClient([]))
    assert out2["index"] == 3 and out2["action"] == "fill"


def test_fill_retries_once_before_degrading(monkeypatch, tmp_path):
    """Этап 3 переработки 2026-08-20: минутный блип модели — не повод для
    заглушки. Первый сбой заполнения тихо ретраится, и только вторая осечка
    подряд деградирует слайд на blank (жалоба №3 — «пайплайн всё ещё ошибается»)."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler
    calls = []

    def _boom_once(client, library, sp, **kw):
        calls.append(sp.index)
        if len(calls) == 1:
            raise filler.FillError("минутный блип")
        return _fake_fill(client, library, sp, **kw)
    monkeypatch.setattr(filler, "fill_slide", _boom_once)
    draft.save_plan("s4c", _outline_plan(question=False))

    out = glass.step_fill("s4c", client=FakeClient([]))
    assert calls == [2, 2]
    assert out["failed"] == [] and out["action"] == "fill"
    s = draft.load_plan("s4c").slides[1]
    assert s.filled and s.status is None and s.template_id == "three-col"


def test_fill_does_not_retry_after_a_slow_failure(monkeypatch, tmp_path):
    """Медленная осечка = таймаут провайдера: второй заход удвоил бы шаг до
    ~10 минут и вылетел бы за клиентские таймауты. Ретраим только быстрые сбои
    (битый JSON, 5xx), медленные деградируют сразу."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler
    calls = []
    ticks = iter([0.0, 300.0])            # первая попытка «шла» 300 секунд
    monkeypatch.setattr(glass, "_now", lambda: next(ticks))

    def _boom(client, library, sp, **kw):
        calls.append(sp.index)
        raise filler.FillError("таймаут")
    monkeypatch.setattr(filler, "fill_slide", _boom)
    draft.save_plan("s4d", _outline_plan(question=False))

    out = glass.step_fill("s4d", client=FakeClient([]))
    assert calls == [2]                   # без второго захода
    assert out["failed"] == [2]


# ── общий клиент шага: липкий резерв переживает шаги ─────────────────────────
def test_glass_client_is_shared_and_preflighted(monkeypatch):
    """Этап 3 переработки 2026-08-20: один клиент на процесс.

    Раньше _kimi() создавал нового клиента на КАЖДЫЙ шаг: липкое переключение
    на резерв (инцидент 2026-08-04) сбрасывалось, и каждый шаг заново платил
    полным таймаутом по мёртвой основной модели. Preflight — один раз при
    создании: мёртвая модель обнаруживается до первого шага, а не внутри него."""
    import htmlslides.pipeline.client as clientmod
    made = []

    class _FakeKimi:
        def __init__(self, **kw):
            made.append(kw)
            self.preflights = 0

        def preflight(self, notice=None):
            self.preflights += 1
    monkeypatch.setattr(clientmod, "KimiClient", _FakeKimi)
    monkeypatch.setattr(glass, "_CLIENT", None)

    a = glass._kimi()
    b = glass._kimi()
    assert a is b and len(made) == 1
    assert a.preflights == 1


def test_glass_client_survives_dead_preflight(monkeypatch):
    """Молчат обе модели — клиент всё равно создаётся: каждый шаг деградирует
    честно сам (failed-слайд, эвристика скоринга), а не 500 на весь конвейер."""
    import htmlslides.pipeline.client as clientmod

    class _DeadKimi:
        def __init__(self, **kw):
            pass

        def preflight(self, notice=None):
            raise RuntimeError("обе модели молчат")
    monkeypatch.setattr(clientmod, "KimiClient", _DeadKimi)
    monkeypatch.setattr(glass, "_CLIENT", None)

    assert glass._kimi() is not None
    assert glass._kimi() is glass._kimi()  # и он остаётся общим


# ── answer ───────────────────────────────────────────────────────────────────
def test_answer_is_instant_and_the_step_fills_the_slide(monkeypatch, tmp_path):
    """Ответ мгновенно записывает выбор в план БЕЗ вызова модели: заполняет
    слайд обычный шаг конвейера. Держать HTTP-ответ на всё время fill_slide
    нельзя — при деградации провайдера (сотни секунд на вызов) клиент обрывал
    ожидание по своему таймауту, врал «проверьте интернет», а поздний
    серверный результат панель уже не подхватывала (живой прогон 2026-08-20)."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler
    seen = {}

    def _spy(client, library, sp, **kw):
        seen["brief"] = sp.content.get("brief", "")
        seen["template_id"] = sp.template_id
        return _fake_fill(client, library, sp, **kw)
    monkeypatch.setattr(filler, "fill_slide", _spy)
    draft.save_plan("s5", _outline_plan())

    out = glass.answer("s5", 2, template_id="timeline", message="покажи даты")
    assert out["open_questions"] == []
    assert seen == {}                       # ответ не тратит ни вызова модели
    s = draft.load_plan("s5").slides[1]
    assert s.template_id == "timeline" and not s.filled
    assert s.status is None and s.question is None and s.candidates is None

    out2 = glass.step_fill("s5", client=FakeClient([]))
    assert out2["action"] == "fill" and out2["index"] == 2
    assert draft.load_plan("s5").slides[1].filled
    assert "Уточнение автора: покажи даты" in seen["brief"]
    assert seen["template_id"] == "timeline"


def test_answer_kind_rides_the_plan_to_the_fill(monkeypatch, tmp_path):
    """Тип схемы из мастера «Схема» доезжает до филлера через план: ответ
    мгновенный, заполняет уже следующий шаг — и он обязан знать выбранную
    автором «воронку», а не угадывать её заново."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler
    seen = {}

    def _spy(client, library, sp, *, deck_title="", extra="", diagram_kind=""):
        seen["kind"] = diagram_kind
        return _fake_fill(client, library, sp, deck_title=deck_title)
    monkeypatch.setattr(filler, "fill_slide", _spy)
    draft.save_plan("s5k", _outline_plan())

    glass.answer("s5k", 2, template_id="diagram", kind="funnel")
    assert draft.load_plan("s5k").slides[1].diagram_kind == "funnel"
    out = glass.step_fill("s5k", client=FakeClient([]))
    assert out["action"] == "fill" and out["index"] == 2
    assert seen["kind"] == "funnel"


def _diagram_fields():
    from webapp import slide_types
    return slide_types.validate_fields("diagram", {
        "heading": "Как проходит заявка", "subtitle": "",
        "diagram": {"kind": "process", "nodes": [
            {"id": "a", "label": "Заявка"}, {"id": "b", "label": "Запуск"}]}})


def test_answer_away_from_a_diagram_drops_its_typed_fields(monkeypatch, tmp_path):
    """Ответ «нет, это не схема» должен менять СЛАЙД, а не только template_id.

    draft_render рисует typed-слайд из slide_type+fields и игнорирует
    template_id — уцелевшие поля прежней схемы оставляли на слайде ту же
    диаграмму: автор выбирал другой макет, платил за перезаполнение и не видел
    никакой разницы."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler
    monkeypatch.setattr(filler, "fill_slide", _fake_fill)
    plan = _outline_plan()
    plan.slides[1] = plan.slides[1].model_copy(update={
        "template_id": "diagram", "slide_type": "diagram",
        "fields": _diagram_fields(), "candidates": ["diagram", "stats-row"]})
    draft.save_plan("s5d", plan)

    glass.answer("s5d", 2, template_id="stats-row")
    glass.step_fill("s5d", client=FakeClient([]))    # заполняет ответивший слайд
    s = draft.load_plan("s5d").slides[1]
    assert s.template_id == "stats-row"
    assert s.slide_type is None and s.fields is None
    from webapp.paths import session_dir
    # не по классу и не по атрибуту: и CSS деки, и инлайн-движок поминают их
    # всегда — ищем сам хост, который ставит шаблон diagram.html
    assert '<div class="diagram-host"' not in (
        session_dir("s5d") / "deck.html").read_text("utf-8")


def test_failed_fill_hides_the_diagram_it_replaced(monkeypatch, tmp_path):
    """Аварийный blank не должен прятаться за прежней схемой."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler

    def _boom(client, library, sp, **kw):
        raise filler.FillError("slide 2: сеть упала")
    monkeypatch.setattr(filler, "fill_slide", _boom)
    plan = _outline_plan()
    plan.slides[1] = plan.slides[1].model_copy(update={
        "template_id": "diagram", "slide_type": "diagram",
        "fields": _diagram_fields()})
    draft.save_plan("s5e", plan)

    glass.answer("s5e", 2, message="перезаполни")
    glass.step_fill("s5e", client=FakeClient([]))
    s = draft.load_plan("s5e").slides[1]
    assert s.template_id == "blank" and s.status == "failed"
    assert s.slide_type is None and s.fields is None


def test_answer_rejects_bad_index_and_unknown_template(monkeypatch, tmp_path):
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    draft.save_plan("s6", _outline_plan())
    with pytest.raises(IndexError):
        glass.answer("s6", 99)
    with pytest.raises(SlotValidationError):
        glass.answer("s6", 2, template_id="ghost")
    # план не испорчен отказом
    assert draft.load_plan("s6").slides[1].status == "needs_input"


# ── refill_slide (смена макета с перезаполнением) ────────────────────────────
def test_refill_uses_the_brief_and_needs_no_question(monkeypatch, tmp_path):
    """Готовый слайд перезаполняется по исходному фрагменту документа.

    Ответ на вопрос (answer) для этого не годится: он требует status
    needs_input/failed, а тут автор правит обычный заполненный слайд."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler
    seen = {}

    def _spy(client, library, sp, **kw):
        seen["brief"] = sp.content.get("brief", "")
        seen["template_id"] = sp.template_id
        return _fake_fill(client, library, sp, **kw)
    monkeypatch.setattr(filler, "fill_slide", _spy)
    plan = _outline_plan(question=False)
    plan.slides[2] = plan.slides[2].model_copy(update={
        "filled": True, "content": {"title": "старое"}})
    draft.save_plan("r1", plan)

    out = glass.refill_slide("r1", 3, template_id="timeline",
                             client=FakeClient([]))
    assert out["index"] == 3
    s = draft.load_plan("r1").slides[2]
    assert s.template_id == "timeline" and s.filled
    assert s.content["title"] == "Слайд 3"      # содержимое написано заново
    assert seen["brief"] == "метрики" and seen["template_id"] == "timeline"
    assert s.brief == "метрики"                 # фрагмент документа цел


def test_refill_without_a_brief_falls_back_to_the_slide_text(monkeypatch,
                                                             tmp_path):
    """Слайд, добавленный руками, брифа не имеет — контекст берём с него самого.

    Иначе действие на таких слайдах молча выдавало бы пустой слайд. В план этот
    текст не попадает: brief в панели подписан «фрагмент исходного документа»."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler
    seen = {}

    def _spy(client, library, sp, **kw):
        seen["brief"] = sp.content.get("brief", "")
        return _fake_fill(client, library, sp, **kw)
    monkeypatch.setattr(filler, "fill_slide", _spy)
    plan = draft.DraftPlan(title="Т", slides=[
        draft.DraftSlide(template_id="timeline", filled=True, content={
            "title": "Что сделали", "items": ["перенесли базу", "выключили старую"],
            "image": "data:image/png;base64,AAA"})])
    draft.save_plan("r2", plan)

    glass.refill_slide("r2", 1, template_id="stats-row", client=FakeClient([]))
    assert "Что сделали" in seen["brief"] and "выключили старую" in seen["brief"]
    assert "data:" not in seen["brief"]         # картинка — не текст
    assert draft.load_plan("r2").slides[0].brief == ""


def test_refill_rejects_bad_index_template_and_empty_slide(monkeypatch, tmp_path):
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    draft.save_plan("r3", draft.DraftPlan(title="Т", slides=[
        draft.DraftSlide(template_id="three-col", brief="тема"),
        draft.DraftSlide(template_id="blank")]))
    with pytest.raises(IndexError):
        glass.refill_slide("r3", 99, client=FakeClient([]))
    with pytest.raises(SlotValidationError):
        glass.refill_slide("r3", 1, template_id="ghost", client=FakeClient([]))
    with pytest.raises(KeyError):
        glass.refill_slide("r3", 1, kind="ghost-kind", client=FakeClient([]))
    # Пустой слайд: заполнять нечем — отказ, а не выдуманное содержимое.
    with pytest.raises(glass.NoContext):
        glass.refill_slide("r3", 2, template_id="stats-row",
                           client=FakeClient([]))


def test_refill_failure_degrades_like_a_step(monkeypatch, tmp_path):
    """Осечка модели не теряет слайд: заглушка с темой + объяснение."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler

    def _boom(client, library, sp, **kw):
        raise filler.FillError("сеть упала")
    monkeypatch.setattr(filler, "fill_slide", _boom)
    plan = _outline_plan(question=False)
    plan.slides[2] = plan.slides[2].model_copy(update={"filled": True})
    draft.save_plan("r4", plan)

    out = glass.refill_slide("r4", 3, template_id="timeline",
                             client=FakeClient([]))
    assert out["failed"] == [3]
    s = draft.load_plan("r4").slides[2]
    assert s.template_id == "blank" and s.status == "failed"
    assert s.content.get("title") and s.question


def test_old_plan_json_valid_without_glass_fields():
    """Старые plan.json без status/question/candidates валидны без миграции."""
    plan = draft.DraftPlan.model_validate_json(
        '{"title":"Т","slides":[{"template_id":"cover","content":{}}]}')
    s = plan.slides[0]
    assert s.status is None and s.question is None and s.candidates is None


# ── endpoints /glass/* ───────────────────────────────────────────────────────
def _new_draft(c):
    return c.post("/api/drafts", json={"mode": "manual"},
                  headers=H()).json()["session_id"]


def test_glass_start_endpoint(monkeypatch, tmp_path):
    with _client_app(monkeypatch, tmp_path) as c:
        sid = _new_draft(c)
        seen = {}

        def _fake_start(session_id, source, *, origin_name=None):
            seen["args"] = (session_id, source.name, source.read_bytes())
            seen["origin"] = origin_name
            plan = draft.DraftPlan(title="Т", slides=[
                draft.DraftSlide(template_id="cover", filled=True)])
            draft.save_plan(session_id, plan)   # как настоящий start_glass
            return plan
        monkeypatch.setattr(glass, "start_glass", _fake_start)

        r = c.post(f"/api/drafts/{sid}/glass/start", headers=H(),
                   files={"file": ("doc.md", "# Т\n\nтекст".encode(), "text/markdown")})
        assert r.status_code == 200
        assert r.json()["title"] == "Т"
        assert seen["args"][0] == sid and seen["args"][1] == "input.md"
        # на диск файл ложится как input.md, но имя автора обязано доехать до
        # сборки — иначе обложка документа без титула выходит с надписью «INPUT»
        assert seen["origin"] == "doc.md"

        # непустой черновик → 409 (стекло стартует с чистой сессии)
        assert c.post(f"/api/drafts/{sid}/glass/start", headers=H(),
                      files={"file": ("doc.md", b"# T", "text/markdown")}
                      ).status_code == 409


def test_glass_start_validation(monkeypatch, tmp_path):
    with _client_app(monkeypatch, tmp_path) as c:
        sid = _new_draft(c)
        r = c.post(f"/api/drafts/{sid}/glass/start", headers=H(),
                   files={"file": ("doc.exe", b"x", "application/x-msdownload")})
        assert r.status_code == 400
        import re
        # A-2: текст ошибки русский, как остальные 400-ки
        assert re.search("[а-яА-Я]", r.json()["detail"]), r.json()
        assert c.post(f"/api/drafts/{sid}/glass/start", headers=H(),
                      files={"file": ("doc.md", b"   \n", "text/markdown")}
                      ).status_code == 400
        # сбой parse → 500 с человеческой причиной, не голый traceback
        def _boom(session_id, source):
            raise ValueError("bad doc")
        monkeypatch.setattr(glass, "start_glass", _boom)
        r = c.post(f"/api/drafts/{sid}/glass/start", headers=H(),
                   files={"file": ("doc.md", b"# T\n\ntext", "text/markdown")})
        assert r.status_code == 500
        assert "не удалось разобрать" in r.json()["detail"]
        # чужая сессия → 404 (изоляция владения)
        assert c.post(f"/api/drafts/{sid}/glass/start", headers=H("intruder"),
                      files={"file": ("doc.md", b"# T\n\ntext", "text/markdown")}
                      ).status_code == 404


def test_glass_step_and_answer_endpoints(monkeypatch, tmp_path):
    with _client_app(monkeypatch, tmp_path) as c:
        sid = _new_draft(c)
        monkeypatch.setattr(glass, "step_fill", lambda session_id: {
            "done": True, "index": None, "action": None,
            "open_questions": [], "plan": {}})
        r = c.post(f"/api/drafts/{sid}/glass/step", headers=H())
        assert r.status_code == 200 and r.json()["done"] is True

        # /glass/score — параллельный разведчик клиентского цикла
        monkeypatch.setattr(glass, "score_next", lambda session_id: {
            "done": False, "index": 2, "action": "score",
            "open_questions": [], "unscored": [3], "plan": {}})
        r = c.post(f"/api/drafts/{sid}/glass/score", headers=H())
        assert r.status_code == 200 and r.json()["action"] == "score"
        assert c.post(f"/api/drafts/{sid}/glass/score", headers=H("intruder")
                      ).status_code == 404

        seen = {}

        def _fake_answer(session_id, index, *, template_id=None, kind=None,
                         message=""):
            if index < 1 or index == 99:      # как настоящий answer
                raise IndexError(f"slide {index} out of range")
            if template_id == "ghost":
                raise SlotValidationError("unknown template: ghost")
            if kind == "ghost-kind":
                raise KeyError("unknown diagram kind: ghost-kind")
            seen["call"] = (session_id, index, template_id, kind, message)
            return {"index": index, "open_questions": [], "plan": {}}
        monkeypatch.setattr(glass, "answer", _fake_answer)

        r = c.post(f"/api/drafts/{sid}/glass/answer", headers=H(),
                   json={"index": 2, "template_id": "timeline",
                         "message": "покажи даты"})
        assert r.status_code == 200
        assert seen["call"] == (sid, 2, "timeline", None, "покажи даты")
        # Тип схемы едет данными: выбор автора в мастере «Схема» — требование к
        # филлеру, а не пожелание в тексте, которое модель вольна не заметить.
        r = c.post(f"/api/drafts/{sid}/glass/answer", headers=H(),
                   json={"index": 2, "template_id": "diagram", "kind": "funnel"})
        assert r.status_code == 200
        assert seen["call"] == (sid, 2, "diagram", "funnel", "")
        assert c.post(f"/api/drafts/{sid}/glass/answer", headers=H(),
                      json={"index": 2, "kind": "ghost-kind"}).status_code == 400
        assert c.post(f"/api/drafts/{sid}/glass/answer", headers=H(),
                      json={"message": "без индекса"}).status_code == 400
        assert c.post(f"/api/drafts/{sid}/glass/answer", headers=H(),
                      json={"index": 99}).status_code == 400
        assert c.post(f"/api/drafts/{sid}/glass/answer", headers=H(),
                      json={"index": 2, "template_id": "ghost"}).status_code == 400


# ── конкурентность: два писателя plan.json вокруг долгого вызова LLM ─────────
# Типовой триггер — не «две вкладки», а F5 в середине 30-секундного шага:
# соединение рвётся, а серверный поток продолжает работать и дописывает свой
# снимок поверх нового. Гейт ниже воспроизводит перехлёст детерминированно.
def _gated_fill(gate, released):
    """fill_slide, который замирает на входе, пока тест не отпустит гейт."""
    def _fill(client, library, sp, *, deck_title="", extra="", diagram_kind=""):
        released.set()
        gate.wait(5)
        return _fake_fill(client, library, sp, deck_title=deck_title)
    return _fill


def test_step_does_not_clobber_answer_written_during_llm_call(monkeypatch, tmp_path):
    """Ответ автора, пришедший во время шага, не откатывается снимком шага."""
    import threading

    import htmlslides.pipeline.filler as filler
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    gate, entered = threading.Event(), threading.Event()
    monkeypatch.setattr(filler, "fill_slide", _gated_fill(gate, entered))
    # Вопрос в хвосте: перед ним шагу есть что заполнять, и работа шага
    # перекрывается с ответом автора — иначе сборка просто стоит на паузе.
    draft.save_plan("c1", _outline_plan(at=4))

    step = threading.Thread(target=glass.step_fill,
                            kwargs={"session_id": "c1", "client": FakeClient([])})
    step.start()
    assert entered.wait(5)                      # шаг взял слайд 2 и завис в LLM
    monkeypatch.setattr(filler, "fill_slide", _fake_fill)
    glass.answer("c1", 4, template_id="timeline",
                 message="покажи этапы")        # автор ответил на вопрос слайда 4
    gate.set()
    step.join(5)

    s4 = draft.load_plan("c1").slides[3]
    assert s4.template_id == "timeline" and s4.status is None and not s4.filled
    assert "Уточнение автора: покажи этапы" in s4.brief
    # Ответ мгновенный, заполняет слайд конвейер шагов: 3 (метрики), затем 4.
    glass.step_fill("c1", client=FakeClient([]))
    glass.step_fill("c1", client=FakeClient([]))
    assert draft.load_plan("c1").slides[3].filled


def test_two_steps_never_take_the_same_slide(monkeypatch, tmp_path):
    """Параллельные шаги берут разные слайды — вызов модели не оплачивается дважды."""
    import threading

    import htmlslides.pipeline.filler as filler
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    seen, guard = [], threading.Lock()
    ready = threading.Barrier(2, timeout=5)

    def _fill(client, library, sp, *, deck_title="", extra="", diagram_kind=""):
        with guard:
            seen.append(sp.index)
        ready.wait()                             # оба шага заняты одновременно
        return _fake_fill(client, library, sp, deck_title=deck_title)

    monkeypatch.setattr(filler, "fill_slide", _fill)
    draft.save_plan("c2", _outline_plan(question=False))

    threads = [threading.Thread(target=glass.step_fill,
                                kwargs={"session_id": "c2", "client": FakeClient([])})
               for _ in range(2)]
    for t in threads:
        t.start()
    for t in threads:
        t.join(5)

    assert sorted(seen) == [2, 3]
    plan = draft.load_plan("c2")
    assert plan.slides[1].filled and plan.slides[2].filled


def test_answer_refuses_a_slide_that_is_not_waiting(monkeypatch, tmp_path):
    """Устаревшая вкладка не перезаполняет обложку и готовые слайды."""
    import htmlslides.pipeline.filler as filler
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    monkeypatch.setattr(filler, "fill_slide", _fake_fill)
    draft.save_plan("c3", _outline_plan())

    with pytest.raises(glass.AnswerNotExpected):
        glass.answer("c3", 1, message="перезаполни")
    assert draft.load_plan("c3").slides[0].content == {"title": "Т"}


def test_failed_fill_after_answer_still_leaves_a_card(monkeypatch, tmp_path):
    """Ответ мгновенный и без модели; сбой на заполнении шага оставляет карточку."""
    import htmlslides.pipeline.filler as filler
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    draft.save_plan("c4", _outline_plan())

    def _boom_kimi(*a, **kw):
        raise AssertionError("ответ не должен трогать модель")
    monkeypatch.setattr(glass, "_kimi", _boom_kimi)

    glass.answer("c4", 2, template_id="timeline", message="этапы")
    s2 = draft.load_plan("c4").slides[1]
    assert s2.status is None and s2.template_id == "timeline"

    def _boom(*a, **kw):
        raise filler.FillError("провайдер недоступен")
    monkeypatch.setattr(filler, "fill_slide", _boom)
    glass.step_fill("c4", client=FakeClient([]))
    s2 = draft.load_plan("c4").slides[1]
    assert s2.status == "failed" and "Не удалось заполнить" in s2.question


def test_answer_over_api_reports_conflict_not_bad_request(monkeypatch, tmp_path):
    import htmlslides.pipeline.filler as filler
    client = _client_app(monkeypatch, tmp_path)
    monkeypatch.setattr(filler, "fill_slide", _fake_fill)
    with client:
        sid = client.post("/api/drafts", headers=H()).json()["session_id"]
        draft.save_plan(sid, _outline_plan())
        r = client.post(f"/api/drafts/{sid}/glass/answer",
                        json={"index": 1, "message": "мимо"}, headers=H())
    assert r.status_code == 409


def test_chat_build_outline_clears_open_glass_questions(monkeypatch, tmp_path):
    """Черновик у степпера и чат-агента общий: автор может уйти из стеклянной
    сборки и дособрать деку кнопкой чата. Заполненный там слайд не должен
    остаться с карточкой «ИИ сомневается в макете» — вопрос уже неактуален."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler
    from webapp import chat_agent, draft_render
    monkeypatch.setattr(filler, "fill_slide", _fake_fill)
    monkeypatch.setattr(chat_agent, "_pick_template", lambda *a, **kw: "three-col")
    monkeypatch.setattr(draft_render, "render_draft", lambda *a, **kw: None)
    draft.save_plan("cb1", _outline_plan())

    chat_agent.build_outline("cb1", client=FakeClient([]))

    plan = draft.load_plan("cb1")
    assert all(s.filled for s in plan.slides)
    assert [s.status for s in plan.slides] == [None] * 4
    assert plan.slides[1].question is None and plan.slides[1].candidates is None


def test_glass_start_refuses_past_the_per_user_ceiling(monkeypatch, tmp_path):
    """Стеклянная сборка идёт мимо очереди раннера, значит и мимо его лимита:
    без своего потолка автор держал бы сколько угодно недозаполненных аутлайнов."""
    monkeypatch.setattr(glass, "MAX_ACTIVE_PER_USER", 2)
    client = _client_app(monkeypatch, tmp_path)
    doc = tmp_path / "d.md"
    doc.write_text("# Т\n\n## Раздел\n\nТекст раздела.\n", encoding="utf-8")
    with client:
        for _ in range(2):                       # два аутлайна «в работе»
            sid = client.post("/api/drafts", headers=H()).json()["session_id"]
            draft.save_plan(sid, _outline_plan())
        sid = client.post("/api/drafts", headers=H()).json()["session_id"]
        r = client.post(f"/api/drafts/{sid}/glass/start", headers=H(),
                        files={"file": ("d.md", doc.read_bytes(), "text/markdown")})
    assert r.status_code == 429


def test_glass_ceiling_not_bypassed_by_fresh_empty_drafts(monkeypatch, tmp_path):
    """B-8 (аудит 2026-08-14): окно подсчёта было limit=MAX+1 НОВЕЙШИХ черновиков
    — достаточно завести несколько свежих пустых, и незавершённые аутлайны
    «выпадали» из окна: лимит на вызовы модели/диск обходился тривиально."""
    monkeypatch.setattr(glass, "MAX_ACTIVE_PER_USER", 2)
    client = _client_app(monkeypatch, tmp_path)
    doc = tmp_path / "d.md"
    doc.write_text("# Т\n\n## Раздел\n\nТекст раздела.\n", encoding="utf-8")
    with client:
        for _ in range(2):                       # два аутлайна «в работе»
            sid = client.post("/api/drafts", headers=H()).json()["session_id"]
            draft.save_plan(sid, _outline_plan())
        for _ in range(4):                       # свежие ПУСТЫЕ вытесняют их из окна
            client.post("/api/drafts", headers=H())
        sid = client.post("/api/drafts", headers=H()).json()["session_id"]
        r = client.post(f"/api/drafts/{sid}/glass/start", headers=H(),
                        files={"file": ("d.md", doc.read_bytes(), "text/markdown")})
    assert r.status_code == 429, r.text


def test_unfinished_outlines_counts_only_pending(monkeypatch, tmp_path):
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    draft.save_plan("u1", _outline_plan())                       # есть незаполненные
    draft.save_plan("u2", draft.DraftPlan(slides=[               # всё заполнено
        draft.DraftSlide(template_id="cover", brief="т", filled=True)]))
    assert glass.unfinished_outlines(["u1", "u2", "нет-такой"]) == 1


def test_start_glass_caps_sections_and_says_so(monkeypatch, tmp_path):
    """Потолок разделов ниже, чем у чёрного ящика: стеклянную деку автор
    дособирает по слайду. Обрезку не молчим — иначе автор гадает, куда делись
    разделы (уведомление живёт в плане: ответ /glass/start до панели не доживает,
    страница уходит в редактор)."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    monkeypatch.setattr(glass, "MAX_GLASS_SLIDES", 3)
    src = tmp_path / "big.md"
    src.write_text("# Док\n\n" + "".join(
        f"## Раздел {i}\n\nТекст раздела {i} про метрики.\n\n" for i in range(5)),
        encoding="utf-8")

    plan = glass.start_glass("cap1", src)

    assert len(plan.slides) == 4              # обложка + 3 раздела
    assert "3" in plan.notice and "2" in plan.notice
    # Хвост назван числом и адресом: с него продолжит вторая дека.
    assert plan.rest == 2 and plan.rest_from == 3


def _capped_deck(monkeypatch, tmp_path, sections=5, cap=3):
    """Обрезанная по потолку дека, исходник которой лежит в сессии (как после
    /glass/start: загрузка всегда ложится в session_dir как `input.<ext>`)."""
    from webapp.paths import session_dir
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    monkeypatch.setattr(glass, "MAX_GLASS_SLIDES", cap)
    src = session_dir("cap1") / "input.md"
    src.parent.mkdir(parents=True, exist_ok=True)
    src.write_text("# Док\n\n" + "".join(
        f"## Раздел {i}\n\nТекст раздела {i} про метрики.\n\n"
        for i in range(sections)), encoding="utf-8")
    return glass.start_glass("cap1", src)


def test_continue_glass_takes_the_tail_into_a_second_deck(monkeypatch, tmp_path):
    """Потерянный хвост возвращается второй декой, а не резкой файла руками.

    Потолок поднимать некуда (сорок слайдов автор и так проходит по одному
    больше получаса), поэтому обратимой делаем саму потерю: исходник уже лежит
    в сессии, номер первого невзятого раздела — в плане."""
    from webapp.paths import session_dir
    first = _capped_deck(monkeypatch, tmp_path)

    second = glass.continue_glass("cap1", "cap2")

    assert len(second.slides) == 3            # обложка + ровно хвост (3 и 4)
    assert "Раздел 3" in second.slides[1].brief
    assert "Раздел 4" in second.slides[2].brief
    assert all(s.status == "unscored" for s in second.slides[1:])   # тоже лениво
    assert second.rest == 0 and second.rest_from == 0   # хвоста больше нет
    assert second.title == first.title        # титул наследуем: на диске «input»
    assert (session_dir("cap2") / "input.md").is_file()  # исходник свой
    back = draft.load_plan("cap1")
    assert back.slides[1].brief == first.slides[1].brief  # первая дека не тронута
    # …кроме снятой заявки: хвост забран, второй раз предлагать нечего.
    assert back.rest == 0 and back.rest_from == 0


def test_continue_glass_refuses_when_nothing_is_left(monkeypatch, tmp_path):
    """Продолжать нечего — говорим об этом, а не заводим пустую вторую деку."""
    _capped_deck(monkeypatch, tmp_path, sections=2, cap=3)
    with pytest.raises(ValueError):
        glass.continue_glass("cap1", "cap2")


def test_glass_rest_endpoint_returns_a_new_draft(monkeypatch, tmp_path):
    """API-контракт кнопки: новая сессия-черновик с хвостом, старая не тронута."""
    monkeypatch.setattr(glass, "MAX_GLASS_SLIDES", 3)
    doc = ("# Док\n\n" + "".join(f"## Раздел {i}\n\nТекст раздела {i}.\n\n"
                                 for i in range(5))).encode()
    with _client_app(monkeypatch, tmp_path) as c:
        sid = _new_draft(c)
        assert c.post(f"/api/drafts/{sid}/glass/start", headers=H(),
                      files={"file": ("Док.md", doc, "text/markdown")}
                      ).status_code == 200

        r = c.post(f"/api/drafts/{sid}/glass/rest", headers=H())
        assert r.status_code == 200
        rest_sid = r.json()["session_id"]
        assert rest_sid != sid and r.json()["slides"] == 2
        assert len(draft.load_plan(rest_sid).slides) == 3
        assert len(draft.load_plan(sid).slides) == 4      # первая дека как была
        # Вторая дека — своя сессия того же автора: она видна в редакторе.
        assert c.get(f"/api/drafts/{rest_sid}", headers=H()).status_code == 200
        assert c.get(f"/api/drafts/{rest_sid}", headers=H("intruder")
                     ).status_code == 404
        # Хвоста больше нет — кнопка не заводит пустых деков ни у второй деки,
        # ни у первой (повторное нажатие с отставшей вкладки).
        assert c.post(f"/api/drafts/{rest_sid}/glass/rest", headers=H()
                      ).status_code == 400
        assert c.post(f"/api/drafts/{sid}/glass/rest", headers=H()
                      ).status_code == 400


def test_step_done_stays_false_while_tail_question_open(monkeypatch, tmp_path):
    """1-4 (аудит раунда 2, agent1): шаг, заполнивший последний «обычный» слайд,
    отвечал done=true при открытом вопросе в хвосте. Семантика: done только
    когда ни вопросов, ни работы не осталось."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler
    monkeypatch.setattr(filler, "fill_slide", _fake_fill)
    draft.save_plan("s14", _outline_plan(at=4))

    out2 = glass.step_fill("s14", client=FakeClient([]))
    assert out2["index"] == 2 and out2["done"] is False
    out3 = glass.step_fill("s14", client=FakeClient([]))
    assert out3["index"] == 3
    assert out3["done"] is False   # вопрос открыт — сборка не «завершена»
    out4 = glass.step_fill("s14", client=FakeClient([]))
    assert out4["action"] is None and out4["open_questions"] == [4]
    assert out4["done"] is False


def test_step_done_stays_false_while_parallel_step_inflight(monkeypatch, tmp_path):
    """1-4 (там же): второй степпер получал done=true, пока первый ещё заполнял
    последний слайд (_next_index не знает про _INFLIGHT). Вторая вкладка видела
    «сборка завершена» на середине заполнения."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    import htmlslides.pipeline.filler as filler
    monkeypatch.setattr(filler, "fill_slide", _fake_fill)
    plan = _outline_plan(question=False)
    for i in (2, 3):    # слайды 2-3 готовы; 4-й «заполняет» параллельный шаг
        plan.slides[i - 1] = plan.slides[i - 1].model_copy(
            update={"filled": True})
    draft.save_plan("s15", plan)
    glass._INFLIGHT["s15"] = {4}
    try:
        out = glass.step_fill("s15", client=FakeClient([]))
        assert out["index"] is None and out["action"] is None
        assert out["done"] is False   # сосед ещё работает
    finally:
        glass._INFLIGHT.pop("s15", None)


def test_answer_bad_index_speaks_russian(monkeypatch, tmp_path):
    """1-6 (аудит раунда 2, agent1): IndexError из glass.answer уходит в 400
    как есть — текст обязан быть русским."""
    monkeypatch.setenv("SLIDESBOT_WORKDIR", str(tmp_path / "sessions"))
    draft.save_plan("s16", _outline_plan())
    with pytest.raises(IndexError) as ei:
        glass.answer("s16", 99)
    assert re.search("[а-яА-Я]", str(ei.value)), str(ei.value)
