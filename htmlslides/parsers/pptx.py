"""PPTX -> InputDoc: слайд = секция; текст-фреймы, таблицы, диаграммы, картинки.

Всё читается python-pptx прямо из XML — рендер (LibreOffice) не нужен. Кроме
верхнеуровневых фигур собираем: текст внутри ГРУПП (рекурсивно), данные ДИАГРАММ
(как таблицу «категория × ряд») и ЗАМЕТКИ докладчика (в Section.notes, не в blocks).
"""
from __future__ import annotations

from pathlib import Path
from typing import Iterator

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import (Block, ImageBlock, InputDoc, ListBlock, Section, TableBlock,
                   TextBlock)


def _walk_shapes(shapes) -> Iterator:
    """Плоский обход фигур: ГРУППЫ разворачиваем рекурсивно (текст внутри группы —
    такой же контент слайда). Порядок — как в XML."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)
        else:
            yield shape


def _fmt_num(value) -> str:
    """Значение ряда диаграммы -> строка: целочисленный float без хвоста «.0»."""
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


def _chart_blocks(shape) -> list[Block]:
    """Диаграмма -> [заголовок?] + таблица «категория × ряды». Любой сбой чтения
    (нестандартный тип графика и т.п.) — молча пропускаем, парсер не должен падать."""
    try:
        chart = shape.chart
        categories = [str(c).strip() for c in chart.plots[0].categories]
        series = list(chart.series)
    except Exception:
        return []
    if not categories or not series:
        return []

    blocks: list[Block] = []
    try:
        if chart.has_title and chart.chart_title.has_text_frame:
            title = chart.chart_title.text_frame.text.strip()
            if title:
                blocks.append(TextBlock(text=title))
    except Exception:
        pass

    def _series_name(s, i: int) -> str:
        try:
            return (s.name or f"Ряд {i}").strip()
        except Exception:
            return f"Ряд {i}"

    values: list[list] = []
    for s in series:
        try:
            values.append(list(s.values))
        except Exception:
            values.append([])

    rows: list[list[str]] = [["Категория"] + [_series_name(s, i)
                                              for i, s in enumerate(series, 1)]]
    for idx, category in enumerate(categories):
        row = [category]
        for col in values:
            row.append(_fmt_num(col[idx] if idx < len(col) else None))
        rows.append(row)
    blocks.append(TableBlock(rows=rows))
    return blocks


# Нестандартные mime SVG (прод 2026-07-29: сторонний экспорт писал `image/svg`)
# нормализуем в честный `image/svg+xml` — иначе data-URI не отрисуется браузером.
_SVG_MIMES = {"image/svg", "image/svg+xml"}


def _picture_block(shape) -> ImageBlock | None:
    """Картинка -> ImageBlock, либо None, если её не прочитать. Парсер не падает.

    `shape.image` работает только для частей, которые python-pptx признал
    ImagePart (реестр стандартных content-type). Прод-кейс 2026-07-29 (job 22/23):
    дека со сторонним экспортом несла 85 картинок с типом `image/svg` — такая
    часть собирается как generic Part без `.image`, и AttributeError ронял всю
    сборку за полсекунды. Фолбэк читает blob напрямую из related part: контент
    сохраняем, а не выбрасываем. None — только если и это не удалось
    (например, картинка-ссылка на внешний файл, у которой blob'а нет вовсе)."""
    try:
        image = shape.image
        return ImageBlock(data=image.blob, mime=image.content_type)
    except Exception:
        pass
    try:
        rId = shape._element.blip_rId
        if not rId:
            return None
        part = shape.part.related_part(rId)
        blob = getattr(part, "blob", None)
        if not blob:
            return None
        mime = getattr(part, "content_type", "") or "application/octet-stream"
        if mime in _SVG_MIMES:
            mime = "image/svg+xml"
        return ImageBlock(data=blob, mime=mime)
    except Exception:
        return None


def _slide_notes(slide) -> str:
    """Заметки докладчика (пусто, если их нет). has_notes_slide — чтобы не плодить
    пустые notes-слайды обращением к геттеру notes_slide."""
    if not slide.has_notes_slide:
        return ""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def parse_pptx(path: str | Path) -> InputDoc:
    prs = Presentation(str(path))
    doc = InputDoc()
    for number, slide in enumerate(prs.slides, start=1):
        title_shape = slide.shapes.title
        title_id = title_shape.shape_id if title_shape is not None else None
        heading = (title_shape.text or "").strip() if title_shape is not None else ""
        if number == 1 and heading and not doc.title:
            doc.title = heading
        section = Section(heading=heading or f"Слайд {number}", level=1)
        for shape in _walk_shapes(slide.shapes):
            if title_id is not None and shape.shape_id == title_id:
                continue
            if getattr(shape, "has_chart", False):
                section.blocks.extend(_chart_blocks(shape))
                continue
            if getattr(shape, "has_table", False):
                section.blocks.append(TableBlock(
                    rows=[[cell.text.strip() for cell in row.cells]
                          for row in shape.table.rows]))
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                block = _picture_block(shape)
                if block is not None:
                    section.blocks.append(block)
                continue
            if shape.has_text_frame:
                paragraphs = [p.text.strip()
                              for p in shape.text_frame.paragraphs if p.text.strip()]
                if len(paragraphs) > 1:
                    section.blocks.append(ListBlock(items=paragraphs))
                elif paragraphs:
                    section.blocks.append(TextBlock(text=paragraphs[0]))
        section.notes = _slide_notes(slide)
        doc.sections.append(section)
    return doc
